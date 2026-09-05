# -*- coding: utf-8 -*-
"""
import-pptx.py — перенос содержимого «SNP Е2Е процесс.pptx» в src/data/snp/process.json
(модель данных — SPEC.md §3, zod-схема — src/data/schema.ts).

Запуск (из корня репозитория):

    python scripts/import-pptx.py

Скрипт детерминирован: повторный прогон на той же презентации даёт побайтово
тот же JSON. Идентификаторы стабильны и не зависят от порядка обхода фигур:
при коллизии базового slug'а суффикс берётся из «номер слайда + shape_id»
(двухфазная генерация, см. IdFactory).

ЭТО ПЕРВАЯ ПОЛОВИНА КОНВЕЙЕРА
-----------------------------
Импорт кладёт в `position` СЫРУЮ геометрию слайда, на которой карточки
накладываются друг на друга (десятки пересекающихся пар). Пригодные к показу
координаты считает вторая половина — `npm run layout` (scripts/layout.ts,
dagre). Порядок обязателен и обратного не имеет.

Чтобы этот порядок не приходилось помнить, есть одна команда:

    npm run data          # import-pptx.py → layout.ts

Исходная геометрия слайда при этом не теряется: она пишется ещё и в
`node.slidePosition` (SPEC §3), и раскладка сидируется именно ею, а не своим
прошлым результатом. Если прогнать только импорт и закоммитить, `npm run check`
покраснеет: tests/mapContract.test.ts сверяет координаты файла с пересчётом.

Структура презентации (проверено на файле «SNP Е2Е процесс.pptx»):
  слайд 1 — титул, данных нет;
  слайд 2 — обзор уровня 1 (контейнеры этапов, боксы групп, боксы внешних систем,
            блоки выходов этапов, связи этап→этап и система→этап);
  слайды 3..6 — детализация этапов 1..4.

Явных связей (stCxn/endCxn) в презентации нет, поэтому принадлежность узла группе
и концы рёбер выводятся геометрически. Всё, что не распозналось однозначно,
не выдумывается, а печатается в отчёте-сверке.

ЧТО В ДОКУМЕНТЕ НЕ ИЗ ПРЕЗЕНТАЦИИ
---------------------------------
Ровно четыре места, и все четыре на виду — потому что правило проекта «не
изобретать процесс» иначе не проверить:

  · OWNER_DECISION_EDGES — рёбра по решению владельца процесса, которых на
    слайде нет (задача process-map-7bz). Отчёт печатает их отдельным блоком;
  · STAGE_INPUT_ENRICHMENT — входы этапа, взятые со слайда ОБЗОРА вместо слайда
    детализации (задача process-map-qjl): текст всё равно из презентации, но
    выбор источника сделан владельцем. Тоже печатается отдельным блоком;
  · STAGE_GROUP_SPLIT — деление узлов этапа на группы, которого на слайде
    детализации нет (задача process-map-028): обзор показывает две группы,
    детализация — один контейнер. Отдельный блок в отчёте;
  · ручные поля `screen`/`owner` — их проставляет человек в редакторе, см.
    ниже.

Направление data-узлов (`node.direction`, задача process-map-24p) к этому
списку НЕ относится: оно читается из презентации так же, как всё прочее, —
по тому, из какой фигуры узел родился (левая колонка входов слайда детализации
против блока выходов этапа на слайде обзора). Из КООРДИНАТ оно не выводится:
блоки выходов этапов 1 и 2 нарисованы левее середины области шагов.

РУЧНЫЕ ПОЛЯ (`screen`, `owner`) ПЕРЕЖИВАЮТ ПЕРЕГЕНЕРАЦИЮ
---------------------------------------------------------
Ссылок на экраны In.Plan в презентации нет: их проставляет человек в редакторе,
и ради них карта вообще встроена в вики. Скрипт собирает документ из презентации
с нуля, поэтому перед записью он читает предыдущий src/data/snp/process.json и
переносит на новые узлы поля, которых в презентации не существует
(PRESERVED_NODE_FIELDS / PRESERVED_STAGE_FIELDS). Сопоставление — по `id`;
id стабильны по построению (см. IdFactory). Всё, что перенести не удалось,
печатается поимённо — молчаливой потери ссылок быть не должно.

Самопроверка переноса (без презентации, только stdlib):

    python scripts/import-pptx.py --self-test
"""

from __future__ import annotations

import json
import math
import re
import sys
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path
from typing import AbstractSet, Iterable, Sequence

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# --------------------------------------------------------------------------------------
# Константы
# --------------------------------------------------------------------------------------

ROOT = Path(__file__).resolve().parent.parent

MAP_VERSION = "1.0.0"

# Идентификатор карты: имя каталога src/data/<id>/ и ключ overrides в
# localStorage (process-map-3wh.4). Проверяется контрактом карты
# (tests/mapContract.test.ts) на совпадение с именем каталога.
MAP_ID = "snp"

# Дата в шапке обзора («Обновлено …»). Решение владельца от 30.08.2026: это
# дата КАРТЫ, а не презентации.
#
# ПОЧЕМУ КОНСТАНТА, А НЕ ВЫЧИСЛЕНИЕ. Соблазн взять dcterms:modified из самой
# презентации (python-pptx отдаёт его как core_properties.modified) выглядит
# правильным, но отвечает на другой вопрос. Содержание карты — функция ДВУХ
# источников: презентации и таблиц решений владельца в этом же файле
# (OWNER_DECISION_EDGES, STAGE_INPUT_ENRICHMENT, STAGE_GROUP_SPLIT и т.п.).
# У первого метка времени внутри есть, у второго её нет и быть не может.
# Проверяемый факт: презентация с 24.08 не менялась ни разу (один коммит), а
# process.json с тех пор менялся пять раз — правками этого файла. То есть
# автоматика показывала бы 24.08 и продолжала врать.
#
# Дата генерации (date.today()) отпадает по критерию приёмки: два прогона
# конвейера обязаны давать побайтово одинаковый файл, а она делает диффом
# каждый новый день.
#
# ЧТО НЕ ДАЁТ ЕЙ ПРОТУХНУТЬ СНОВА — соседняя константа: отпечаток содержания
# process.json. Расходятся — краснеет tests/updatedAt.test.ts. Проверка на
# TypeScript, а не здесь, потому что Python в CI не запускается вовсе.
MAP_UPDATED_AT = "2026-09-01"

# sha256 содержания src/data/snp/process.json с НЕЙТРАЛИЗОВАННЫМ updatedAt.
#
# Нейтрализация обязательна: иначе обновление даты меняет содержание, значит и
# хеш, значит нужен новый хеш — цикл не сходится. С плейсхолдером правка даты
# отпечатка не требует, а правка данных требует.
#
# ОБНОВЛЯТЬ ВРУЧНУЮ, вместе с MAP_UPDATED_AT. Соблазн «пусть импортёр сам
# пересчитывает при записи» надо отвергнуть: тогда сторож всегда зелёный и
# дата протухнет ровно тем же способом. Python и не смог бы посчитать финальный
# отпечаток — после него файл переписывает scripts/layout.ts.
# Значение печатает падающий тест; алгоритм — в tests/updatedAt.test.ts.
MAP_DATA_FINGERPRINT = "415efc85cc982ddf8cd9c6b419cf8e6ea70cb347641e92b400c77bc7315a51c8"

# Заголовок шапки обзора. Решение владельца от 31.08.2026 (process-map-4d2):
# формулировка макета A1, она же в заголовке PRD.md. Прежняя — «E2E процесс
# планирования In.Plan» — расходилась и с макетом, и с PRD.
#
# Это КОНСТАНТА, а не текст из презентации: правка одного process.json была бы
# затёрта первым же `npm run data`. Задача 4d2 предполагала обратное.
MAP_TITLE = "E2E-процесс планирования поставок"

# Подпись рамки вокруг всего потока этапов на обзоре. Жила в i18n
# (ru.overview.laneFlow), но меняется от карты к карте — значит это свойство
# документа, а не строка интерфейса (process-map-3wh.4).
MAP_MODULE_LABEL = "Модуль SNP"

# --- карта MRP (process-map-3wh.9) -------------------------------------------
# Собирается с ОДНОГО слайда 8 «MRP процесс» (решение владельца: остальные 37
# слайдов вебинара не трогать). Имена констант плоские с суффиксом, а не
# словарь: tests/snp/updatedAt.test.ts читает их регуляркой ^ИМЯ = "значение".
MAP_ID_MRP = "mrp"

# Решение владельца от 01.09.2026: та же расшифровка, что у кода MRP в словаре
# систем (src/i18n/ru.ts) — карта и бейдж называют модуль одинаково.
MAP_TITLE_MRP = "Процесс планирования потребности в материалах"
MAP_MODULE_LABEL_MRP = "Модуль MRP"
MAP_UPDATED_AT_MRP = "2026-09-01"
MAP_DATA_FINGERPRINT_MRP = "394e6ee9b381b0fd01eda89ffc7391993810474ae7c0cdee669023bf429fc9cd"


@dataclass(frozen=True)
class MapSpec:
    """
    Всё, чем одна карта отличается от другой (process-map-3wh.7).

    ПОЧЕМУ РЕЕСТР В КОДЕ, А НЕ КОНФИГ-ФАЙЛ НА ДИСКЕ. Тесты проекта уже разбирают
    ИСХОДНИК этого файла регулярками (tests/snp/importPreserve.test.ts,
    tests/snp/updatedAt.test.ts) — Python в CI не запускается вовсе. Внешний
    конфиг потребовал бы третьего механизма сверки и ещё одного артефакта,
    который надо держать в синхроне.

    ПОЧЕМУ КОНСТАНТЫ MAP_* ОСТАЛИСЬ ПЛОСКИМИ, а не переехали сюда значениями.
    tests/snp/updatedAt.test.ts читает их регуляркой вида ^ИМЯ = "значение".
    Словарь эту регулярку ломает и заставляет переписывать сторож даты; плоские
    имена с суффиксом карты позволяют его параметризовать.
    """

    key: str
    pptx: Path
    json: Path
    required_nodes: Path
    # 'overview+details' — обзор на слайде 2 плюс четыре слайда детализации
    # (устройство презентации SNP). Профиль 'single-slide' вводит process-map-3wh.9.
    profile: str
    slides: int
    # Индекс рабочего слайда для профиля 'single-slide' (0-based). У профиля
    # 'overview+details' слайды заданы устройством презентации: обзор — второй,
    # детализация — с третьего по шестой.
    slide_index: int | None
    map_id: str
    title: str
    module_label: str
    updated_at: str
    fingerprint: str


MAPS: dict[str, MapSpec] = {
    "snp": MapSpec(
        key="snp",
        pptx=ROOT / "SNP Е2Е процесс.pptx",
        json=ROOT / "src" / "data" / "snp" / "process.json",
        required_nodes=ROOT / "tests" / "fixtures" / "snp" / "required-nodes.json",
        profile="overview+details",
        slides=6,
        slide_index=None,
        map_id=MAP_ID,
        title=MAP_TITLE,
        module_label=MAP_MODULE_LABEL,
        updated_at=MAP_UPDATED_AT,
        fingerprint=MAP_DATA_FINGERPRINT,
    ),
    "mrp": MapSpec(
        key="mrp",
        pptx=ROOT / "In.Plan MRP 17-08-2026.pptx",
        json=ROOT / "src" / "data" / "mrp" / "process.json",
        required_nodes=ROOT / "tests" / "fixtures" / "mrp" / "required-nodes.json",
        profile="single-slide",
        slides=38,
        slide_index=7,  # слайд 8 «MRP процесс»
        map_id=MAP_ID_MRP,
        title=MAP_TITLE_MRP,
        module_label=MAP_MODULE_LABEL_MRP,
        updated_at=MAP_UPDATED_AT_MRP,
        fingerprint=MAP_DATA_FINGERPRINT_MRP,
    ),
}

DEFAULT_MAP = "snp"

# 1 px = 9525 EMU (96 dpi). Слайд 12192000 EMU = 1280 px по ширине.
EMU_PER_PX = 9525

SLIDE_WIDTH_EMU = 12192000

# Геометрические пороги (EMU).
CONTAINER_MIN_WIDTH = 2_000_000      # контейнер группы — широкий бесфонный прямоугольник
GROUP_TITLE_MAX_GAP = 1_200_000      # заголовок группы лежит не выше этого над контейнером
CAPTION_MAX_GAP = 400_000            # подпись-выход лежит вплотную под узлом
CAPTION_MIN_OVERLAP = 0.6            # доля ширины подписи, перекрытая узлом
CAPTION_MIN_GAP = -100_000           # допускаем лёгкое налезание подписи на узел
CAPTION_MIN_NODE_HEIGHT = 250_000    # тонкие полосы подписываются сверху, а не снизу
EDGE_SNAP_DETAIL = 300_000           # конец линии «прилипает» к узлу (слайды детализации)
EDGE_SNAP_SECOND = 900_000           # второй проход для линий с одним разрешённым концом
EDGE_SECOND_RATIO = 1.8              # следующий кандидат должен быть настолько же дальше
EDGE_SNAP_OVERVIEW = 500_000         # то же для обзора: там геометрия заметно свободнее
PROMOTE_MIN_ENDPOINTS = 2            # столько концов линий должно упираться в текстбокс
PROMOTE_SNAP = 50_000                # и упираться вплотную: подпись рядом со стрелкой — не узел
LOOSE_DESC_MAX_DIST = 2_500_000      # предел для «свободной» подписи внутри группы
LEFT_MARGIN_LIMIT = SLIDE_WIDTH_EMU * 15 // 100  # левое поле — колонка входов
KEY_OUTPUT_TOP_OFFSET = 500_000      # блоки выходов — в нижней части контейнера этапа
KEY_OUTPUT_BOTTOM_OFFSET = 700_000
DECOR_ARROW_MAX_WIDTH = 400_000      # мелкие стрелки-коннекторы между боксами обзора

MAX_KEY_OUTPUTS = 4                  # ограничение zod-схемы
STAGE_COUNT = 4                      # ограничение zod-схемы: number ∈ {1,2,3,4}

# Заливка плашек-артефактов (входы и выход процесса) в профиле «одиночный слайд».
# В презентации SNP входы нарисованы надписями, здесь — автофигурами, и без
# отдельного признака они стали бы обычными шагами.
ARTIFACT_FILL = "scheme:accent2"
MAX_ID_LENGTH = 72                   # длиннее, чтобы различающая часть текста не срезалась

SYSTEM_CODES = ("DP", "PS", "IO", "ERP", "MRP", "INPLAN", "BI", "EPM")

# --------------------------------------------------------------------------------------
# Контракт с src/data/schema.ts
# --------------------------------------------------------------------------------------
#
# NODE_KEY_ORDER / STAGE_KEY_ORDER повторяют ПОРЯДОК ключей zod-схем
# ProcessNodeSchema и StageSchema. Это не косметика: экспорт из приложения
# (src/utils/processTransfer.ts::serializeProcessMap) прогоняет карту через zod,
# который пересобирает объекты в порядке схемы, и обязан совпадать с этим файлом
# побайтово. Расхождение ловит tests/importPreserve.test.ts.
NODE_KEY_ORDER = (
    "id",
    "type",
    # Уточнение типа для узлов BPMN. Импортёр презентаций их не пишет, но
    # порядок обязан совпадать со схемой ключ в ключ:
    # tests/snp/importPreserve.test.ts сравнивает этот кортеж со
    # ProcessNodeSchema.shape точно и по порядку.
    "gatewayKind",
    "eventKind",
    "eventDefinition",
    "label",
    "description",
    "group",
    "direction",
    "inputs",
    "outputs",
    "system",
    "owner",
    "screen",
    "position",
    "slidePosition",
)
STAGE_KEY_ORDER = (
    "id",
    "number",
    "title",
    "shortTitle",
    "keyOutputs",
    "warningsCount",
    "screen",
    "groups",
    "nodes",
    "edges",
    "inputs",
    "outputs",
)

# Поля модели, которых В ПРЕЗЕНТАЦИИ НЕТ: их заполняет человек (редактор ссылок —
# SPEC §4.4, `owner` — правкой файла). Импортёр их не создаёт, значит обязан
# переносить из предыдущего process.json, иначе перегенерация их стирает.
PRESERVED_NODE_FIELDS = ("owner", "screen")
PRESERVED_STAGE_FIELDS = ("screen",)

# Остальное импортёр строит сам из презентации.
IMPORTER_NODE_FIELDS = tuple(k for k in NODE_KEY_ORDER if k not in PRESERVED_NODE_FIELDS)
IMPORTER_STAGE_FIELDS = tuple(k for k in STAGE_KEY_ORDER if k not in PRESERVED_STAGE_FIELDS)

# Код возврата, когда ссылки на экраны потеряны: узла с таким id в презентации
# больше нет. Файл при этом записывается — импорт корректен, потерян только
# ручной слой, и его полный список напечатан в отчёте.
EXIT_LINKS_LOST = 2

# --------------------------------------------------------------------------------------
# РЁБРА ПО РЕШЕНИЮ ВЛАДЕЛЬЦА ПРОЦЕССА — ИХ НЕТ В ПРЕЗЕНТАЦИИ
# --------------------------------------------------------------------------------------
#
# ЧИТАТЬ ЦЕЛИКОМ, ПРЕЖДЕ ЧЕМ ДОБАВЛЯТЬ СЮДА СТРОКУ.
#
# Всё остальное в этом файле — чтение слайда: узел есть, потому что на слайде
# есть фигура; ребро есть, потому что на слайде есть линия. Этот список —
# ЕДИНСТВЕННОЕ исключение, и заведено оно ровно для того, чтобы исключение было
# видно. В `stage["edges"]` рёбра отсюда лежат вперемешку с прочитанными со
# слайда и внешне от них не отличаются; отличить их можно только здесь и по
# отчёту импортёра (блок «РЁБРА ПО РЕШЕНИЮ ВЛАДЕЛЬЦА ПРОЦЕССА»).
#
# Условие для новой строки — ЗАФИКСИРОВАННОЕ решение владельца процесса
# (номер задачи обязателен). Догадка «тут по смыслу должна быть стрелка» таким
# решением НЕ является: автоматическое правило «связать узлы группы публикации
# с расчётом» уже предлагалось (задача np4) и было откачено как выдумывание
# процесса. Разница между np4 и тем, что ниже, — не в содержании рёбер, а
# исключительно в источнике, и через полгода восстановить её будет неоткуда,
# кроме этого комментария.
#
# ПОЧЕМУ ПРАВИЛОМ В ИМПОРТЁРЕ, А НЕ ПРАВКОЙ process.json. Импортёр собирает
# документ с нуля, поэтому дописанное руками в JSON стирается следующим же
# прогоном `npm run data` — этот дефект уже чинила задача process-map-2dj для
# ссылок на экраны. Механизм переноса ручных полей (PRESERVED_NODE_FIELDS)
# здесь не подходит: он сопоставляет ПОЛЯ СУЩЕСТВУЮЩИХ узлов по id, а не
# добавляет новые сущности, и «перенос» рёбер означал бы, что импортёр молча
# тянет из старого файла связи, которых в презентации нет, — то есть ровно то,
# что запрещено. Объявление в коде переживает перегенерацию по построению и
# при этом остаётся на виду.
#
# `targets` перечисляет решение ЦЕЛИКОМ, включая концы, которые на слайде уже
# нарисованы: применяется только недостающее, а совпавшее печатается в отчёте
# как пришедшее из презентации. Так строка остаётся читаемой как формулировка
# решения («связать все четыре»), а не как дельта к текущему состоянию слайда,
# и правка презентации не превращает список в тихую ложь.
OWNER_DECISION_EDGES: tuple[dict, ...] = (
    {
        "task": "process-map-7bz",
        "stage": 3,
        "source": "raschet-ogranichennyh-planov",
        "targets": (
            "peredacha-ogranichennogo-prognoza-v-dp",
            "publikaciya-planovyh-zakazov",
            "publikaciya-zayavok-na-zakupku",
            "publikaciya-zayavok-na-peremeschenie",
        ),
        "kind": "process",
        "why": (
            "решение владельца процесса: все четыре узла группы «Публикация планов» "
            "выполняются по результату расчёта ограниченных планов. В презентации "
            "стрелка нарисована только к «Публикация плановых заказов» (линия [146]), "
            "остальные три узла остаются изолированными"
        ),
    },
)

STAGE_INPUT_ENRICHMENT: tuple[dict, ...] = (
    {
        "task": "process-map-qjl",
        "stage": 4,
        # sid исходной фигуры слайда 2 — уходит в IdFactory как происхождение
        # добавленных узлов, чтобы id разводились так же, как у всех прочих.
        "sid": 164,
        "source": (
            "слайд 2, текстбокс [164] — мастер-данные этапа 4; привязан линией [181] "
            "к боксу [8] «Расчёт плана пополнения (DRP/Deployment)»"
        ),
        # Строка есть на слайде обзора и отсутствует на слайде детализации.
        "add": ("Транспортные отношения",),
        # Слева — формулировка слайда 6 (она попадает в модель), справа — слайда 2.
        "expand": (
            ("ОСГ", "Остаточный срок годности (ОСГ)"),
            ("Аллокация", "Резервы (аллокация)"),
            ("Размещенные заказы (PO)", "Размещенные заказы на закупку (PO)"),
        ),
        "why": (
            "решение владельца процесса: слайд обзора уточняет слайд детализации. "
            "Из десяти строк [164] семь совпадают с колонкой входов слайда 6, три "
            "записаны там короче, а «Транспортные отношения» отсутствуют в "
            "презентации где-либо ещё"
        ),
    },
)

OWNER_DECISION_EXTERNAL_IO: tuple[dict, ...] = (
    {
        "task": "process-map-vjz.5",
        "stage": 1,
        "system": "ERP",
        "label": "Управление транзакционными данными",
        "direction": "in",
        "source": "слайд 2, текстбокс [130] — верхний ряд, между [40] и [129]",
        "why": (
            "решение владельца процесса: это самостоятельный источник данных этапа 1, "
            "а не заголовок к соседнему тексту. Автоматика его не берёт, и не может: "
            "ExternalIO строится, только когда сработали И detect_system, И "
            "detect_direction, а во фразе нет ни кода системы, ни направления "
            "(«из модуля» / «в модуль»). Код ERP назван владельцем 31.08.2026 — в "
            "презентации его нет нигде. Ни одна из 13 линий слайда 2 к [130] не "
            "подходит: ближайший конец линии [138] лежит в 352422 EMU ниже и внутри "
            "текстбокса заголовка этапа [15]"
        ),
    },
)

STAGE_GROUP_SPLIT: tuple[dict, ...] = (
    {
        "task": "process-map-028",
        "stage": 4,
        "label": "TLB",
        "nodes": (
            "Формирование рекомендаций с учетом параметров нагрузки на транспорт",
            "Корректировки рекомендаций (заявок на перемещение)",
            "Публикация рекомендаций (заявок на перемещение)",
        ),
        "why": (
            "решение владельца процесса. На слайде 2 у этапа 4 две группы — "
            "«Расчёт плана пополнения (DRP/Deployment)» и «TLB», — а на слайде 6 "
            "контейнер один, и геометрия деления не показывает: равные зазоры "
            "(407049 EMU, разброс 1 EMU), сплошная цепочка стрелок, единый "
            "заголовок по центру над всеми четырьмя шагами. Связка нашлась через "
            "выходы: подпись под первым шагом — «План пополнения по дням», это "
            "выход DRP на обзоре; подписи под вторым и четвёртым — «Схема "
            "загрузки транспортных средств», это выход TLB; третий зажат между "
            "ними в той же цепочке. Заголовок слайда 6 сам делит этап надвое: "
            "«Расчёт плана пополнения (DRP/Deployment) + Транспортбилдер»"
        ),
    },
)

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

TRANSLIT = {
    "а": "a", "б": "b", "в": "v", "г": "g", "д": "d", "е": "e", "ё": "e",
    "ж": "zh", "з": "z", "и": "i", "й": "y", "к": "k", "л": "l", "м": "m",
    "н": "n", "о": "o", "п": "p", "р": "r", "с": "s", "т": "t", "у": "u",
    "ф": "f", "х": "h", "ц": "c", "ч": "ch", "ш": "sh", "щ": "sch",
    "ъ": "", "ы": "y", "ь": "", "э": "e", "ю": "yu", "я": "ya",
}


# --------------------------------------------------------------------------------------
# Вспомогательные структуры
# --------------------------------------------------------------------------------------


@dataclass
class Box:
    """Прямоугольник фигуры в EMU."""

    left: int
    top: int
    width: int
    height: int

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height

    @property
    def cx(self) -> float:
        return self.left + self.width / 2

    @property
    def cy(self) -> float:
        return self.top + self.height / 2

    @property
    def area(self) -> int:
        return self.width * self.height

    def contains_point(self, x: float, y: float) -> bool:
        return self.left <= x <= self.right and self.top <= y <= self.bottom

    def distance_to_point(self, x: float, y: float) -> float:
        dx = max(self.left - x, 0.0, x - self.right)
        dy = max(self.top - y, 0.0, y - self.bottom)
        return math.hypot(dx, dy)


@dataclass
class Shape:
    """Нормализованное представление фигуры слайда."""

    sid: int
    kind: str            # 'auto' | 'line' | 'textbox' | 'placeholder' | 'other'
    box: Box
    paragraphs: list[str]
    fill: str | None     # 'srgb:RRGGBB' | 'scheme:name' | 'noFill' | None
    flip_h: bool
    flip_v: bool
    # Поворот фигуры вокруг центра bbox, в градусах (a:xfrm/@rot делится на 60000).
    # У повёрнутого коннектора left/top/width/height описывают НЕПОВЁРНУТУЮ
    # рамку, поэтому концы линии надо доворачивать (process-map-3wh.18).
    rot: float
    head_arrow: bool
    tail_arrow: bool
    # Явные привязки коннектора к фигурам (a:stCxn/a:endCxn).
    #
    # ВНИМАНИЕ, РАСПРОСТРАНЁННОЕ ЗАБЛУЖДЕНИЕ: считалось, что в презентации SNP
    # привязок нет и потому связи приходится выводить геометрически. Это НЕВЕРНО
    # (проверено перебором, process-map-3wh.8): из 66 линий 47 имеют ОБЕ
    # привязки и только 10 не имеют ни одной. Импортёр их просто никогда не
    # читал. Разбор SNP остаётся геометрическим — менять содержание
    # опубликованной карты эта задача не может, — но см. process-map-3wh.16.
    #
    # Профилю «одиночный слайд» привязки нужны: у линии [112] слайда 8 bbox
    # уходит на 1.28 млн EMU ниже слайда (растянутый коннектор после
    # перемещения фигур), геометрия её теряет, а привязка — нет.
    start_sid: int | None = None
    end_sid: int | None = None
    consumed_by: str | None = None   # для отчёта: кто «съел» текстовую фигуру

    @property
    def text(self) -> str:
        return " ".join(self.paragraphs).strip()

    @property
    def has_text(self) -> bool:
        return bool(self.paragraphs)

    def sort_key(self) -> tuple[int, int, int]:
        return (self.box.top, self.box.left, self.sid)


@dataclass
class NodeDraft:
    node_id: str
    node_type: str
    label: str
    box: Box
    group: str | None = None
    # Колонка data-узла на экране детализации (SPEC §3/§4.2, задача
    # process-map-24p): 'in' — вход этапа, 'out' — выход.
    #
    # НЕ ЭВРИСТИКА И НЕ ГЕОМЕТРИЯ, а происхождение фигуры: у импортёра ровно
    # два места, где рождается data-узел, и каждое знает направление точно —
    #   · build_stage, шаг 4 «левая колонка входов» слайда детализации → 'in';
    #   · build_process_map, блоки выходов этапа под его контейнером на слайде
    #     обзора (слайд 2)                                             → 'out'.
    # Больше никакой путь data-узел не создаёт: node_type_for() возвращает
    # только step/integration/warning. Поэтому поле проставлено у всех
    # data-узлов, а у остальных типов его нет — там оно бессмысленно.
    #
    # Выводить направление из координат нельзя: блоки выходов этапов 1 и 2
    # презентация рисует под контейнером этапа на слайде обзора, и по абсциссе
    # они оказываются ЛЕВЕЕ середины области шагов — прежнее правило «левее
    # середины = вход» давало у этих этапов ноль выходов.
    direction: str | None = None
    description_parts: list[str] = field(default_factory=list)
    inputs: list[str] = field(default_factory=list)
    outputs: list[str] = field(default_factory=list)
    system: str | None = None

    @property
    def description(self) -> str | None:
        return "\n".join(self.description_parts) if self.description_parts else None


@dataclass
class SlideReport:
    slide_no: int
    nodes: int = 0
    data_nodes: int = 0
    groups: int = 0
    edges: int = 0
    lines_total: int = 0
    promoted: list[str] = field(default_factory=list)
    lines_skipped: list[str] = field(default_factory=list)
    text_skipped: list[str] = field(default_factory=list)
    loose_attachments: list[str] = field(default_factory=list)
    questions: list[str] = field(default_factory=list)
    # Артефакты, названные выходом этапа в обзоре, но уже существующие узлом
    # колонки входов на слайде детализации (задача process-map-24p).
    dedup_key_outputs: list[str] = field(default_factory=list)
    # Рёбра, добавленные не из презентации, а по решению владельца процесса
    # (OWNER_DECISION_EDGES, задача process-map-7bz).
    owner_edges: list[str] = field(default_factory=list)
    # Входы, добавленные и переформулированные по слайду обзора
    # (STAGE_INPUT_ENRICHMENT, задача process-map-qjl).
    owner_inputs: list[str] = field(default_factory=list)
    # Внешние системы, названные владельцем, а не выведенные из текста
    # (process-map-vjz.5). Живёт на отчёте слайда 2: таблица привязана к обзору.
    owner_external_io: list[str] = field(default_factory=list)
    # Узлы, повышенные из step в integration по коду системы (process-map-7v1).
    promoted_integrations: list[str] = field(default_factory=list)
    # Узлы, перенесённые в отдельную группу по решению владельца
    # (STAGE_GROUP_SPLIT, задача process-map-028).
    owner_groups: list[str] = field(default_factory=list)
    # Рёбра, у которых хотя бы один конец разрешён явной привязкой коннектора,
    # а не геометрией (process-map-3wh.16).
    cxn_edges: list[str] = field(default_factory=list)


# --------------------------------------------------------------------------------------
# Чтение презентации
# --------------------------------------------------------------------------------------


def normalize_text(value: str) -> str:
    """\x0b — мягкий перенос строки внутри абзаца PowerPoint; схлопываем пробелы."""
    value = value.replace("\x0b", " ").replace(" ", " ")
    return re.sub(r"\s+", " ", value).strip()


def read_paragraphs(shape) -> list[str]:
    """
    Абзацы текстовой фигуры.

    PowerPoint иногда разбивает одну фразу на два абзаца («Информация по будущим» /
    «заказам клиентов»). Признак продолжения — следующий абзац начинается со строчной
    буквы или со служебного символа продолжения; такие абзацы склеиваем.
    Списки (каждый пункт с заглавной буквы) при этом не страдают.
    """
    if not shape.has_text_frame:
        return []
    raw = [normalize_text(p.text) for p in shape.text_frame.paragraphs]
    raw = [p for p in raw if p]
    merged: list[str] = []
    for para in raw:
        first = para[0]
        if merged and (first.islower() or first in "/,;)"):
            merged[-1] = f"{merged[-1]} {para}".strip()
        else:
            merged.append(para)
    return merged


def read_fill(element) -> str | None:
    """Явная заливка из spPr (заливка по теме через p:style здесь не учитывается)."""
    sp_pr = element.find(P + "spPr")
    if sp_pr is None:
        return None
    for child in sp_pr:
        tag = etree.QName(child).localname
        if tag == "noFill":
            return "noFill"
        if tag == "solidFill":
            if len(child) == 0:
                return "solid"
            color = child[0]
            local = etree.QName(color).localname
            if local == "srgbClr":
                return "srgb:" + str(color.get("val")).upper()
            return "scheme:" + str(color.get("val"))
        if tag.endswith("Fill"):
            return tag
    return None


def read_line_ends(element) -> tuple[bool, bool]:
    """Наличие стрелок на концах линии (a:ln/a:headEnd, a:tailEnd)."""
    sp_pr = element.find(P + "spPr")
    if sp_pr is None:
        return (False, False)
    ln = sp_pr.find(A + "ln")
    if ln is None:
        return (False, False)

    def is_arrow(tag: str) -> bool:
        node = ln.find(A + tag)
        if node is None:
            return False
        return str(node.get("type") or "none") != "none"

    return (is_arrow("headEnd"), is_arrow("tailEnd"))


def read_connection(element) -> tuple[int | None, int | None]:
    """
    Привязки коннектора к фигурам: (id начала, id конца) из a:stCxn / a:endCxn.

    Живут в p:cxnSp/p:nvCxnSpPr/p:cNvCxnSpPr. У обычной фигуры этого узла нет,
    у неприкреплённого коннектора — нет атрибутов; и то и другое даёт None.
    """
    nv = element.find(P + "nvCxnSpPr")
    if nv is None:
        return (None, None)
    props = nv.find(P + "cNvCxnSpPr")
    if props is None:
        return (None, None)

    def sid_of(tag: str) -> int | None:
        node = props.find(A + tag)
        if node is None:
            return None
        raw = node.get("id")
        return int(raw) if raw is not None and raw.isdigit() else None

    return (sid_of("stCxn"), sid_of("endCxn"))


def read_rotation(element) -> float:
    """Поворот фигуры в градусах из a:xfrm/@rot (единица — 1/60000 градуса)."""
    sp_pr = element.find(P + "spPr")
    if sp_pr is None:
        return 0.0
    node = sp_pr.find(A + "xfrm")
    if node is None:
        return 0.0
    raw = node.get("rot")
    if raw is None:
        return 0.0
    try:
        return int(raw) / 60000.0
    except ValueError:
        return 0.0


def read_flips(element) -> tuple[bool, bool]:
    sp_pr = element.find(P + "spPr")
    if sp_pr is None:
        return (False, False)
    node = sp_pr.find(A + "xfrm")
    if node is None:
        return (False, False)
    return (node.get("flipH") == "1", node.get("flipV") == "1")


def classify_shape_kind(shape) -> str:
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.LINE:
        return "line"
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        return "textbox"
    if st == MSO_SHAPE_TYPE.PLACEHOLDER:
        return "placeholder"
    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return "auto"
    return "other"


def read_slide(slide) -> list[Shape]:
    shapes: list[Shape] = []
    for shape in slide.shapes:
        if shape.left is None or shape.top is None:
            continue
        element = shape._element  # noqa: SLF001 — python-pptx не даёт публичного доступа к XML
        flip_h, flip_v = read_flips(element)
        head_arrow, tail_arrow = read_line_ends(element)
        start_sid, end_sid = read_connection(element)
        rot = read_rotation(element)
        shapes.append(
            Shape(
                sid=int(shape.shape_id),
                kind=classify_shape_kind(shape),
                box=Box(int(shape.left), int(shape.top), int(shape.width), int(shape.height)),
                paragraphs=read_paragraphs(shape),
                fill=read_fill(element),
                flip_h=flip_h,
                flip_v=flip_v,
                rot=rot,
                head_arrow=head_arrow,
                tail_arrow=tail_arrow,
                start_sid=start_sid,
                end_sid=end_sid,
            )
        )
    shapes.sort(key=Shape.sort_key)
    return shapes


# --------------------------------------------------------------------------------------
# Генерация идентификаторов
# --------------------------------------------------------------------------------------


def transliterate(value: str) -> str:
    return "".join(TRANSLIT.get(ch, ch) for ch in value.lower())


def slugify(value: str, max_length: int = MAX_ID_LENGTH) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", transliterate(value)).strip("-")
    if len(slug) > max_length:
        cut = slug[:max_length]
        if "-" in cut:
            cut = cut[: cut.rindex("-")]
        slug = cut.strip("-")
    return slug or "node"


class IdFactory:
    """
    Стабильные kebab-case id.

    Работает в двух фазах. В первой (collisions=None) фабрика только считает,
    сколько раз встретился каждый базовый slug. Во второй ей передают эту карту:
    уникальный slug становится id как есть, а ВСЕ участники коллизии получают
    суффикс «-<слайд>-<shape_id>». Благодаря этому id не зависит от порядка
    обхода фигур и не меняется, если удалить или переставить соседнюю фигуру.
    """

    def __init__(self, collisions: dict[str, int] | None = None, max_length: int = MAX_ID_LENGTH) -> None:
        self._collisions = collisions
        self._counts: Counter[str] = Counter()
        self._used: set[str] = set()
        self._max_length = max_length

    @property
    def counts(self) -> Counter[str]:
        return self._counts

    def make(self, source: str, slide: int, sid: int, index: int = 0) -> str:
        base = slugify(source, self._max_length)
        self._counts[base] += 1
        if self._collisions is None:
            candidate = f"{base}~{self._counts[base]}"
        else:
            candidate = base if self._collisions.get(base, 0) <= 1 else f"{base}-{slide}-{sid}"
            if candidate in self._used:
                candidate = f"{candidate}-{index}"
            extra = 2
            while candidate in self._used:
                candidate = f"{base}-{slide}-{sid}-{index}-{extra}"
                extra += 1
        self._used.add(candidate)
        return candidate


# --------------------------------------------------------------------------------------
# Правила классификации (SPEC §3: NodeType = step | data | integration | warning)
# --------------------------------------------------------------------------------------

INTEGRATION_FILL = "srgb:A6A6A6"
INTEGRATION_TEXT_RE = re.compile(r"^Передача\b.*\b(из|в)\s+модул", re.IGNORECASE)
WARNING_TEXT_RE = re.compile(r"предупрежден", re.IGNORECASE)
WARNING_ITEM_RE = re.compile(r"^Предупреждени", re.IGNORECASE)
# Как система названа в тексте презентации -> её код в схеме.
#
# ПОЧЕМУ АЛИАСЫ, А НЕ НОВЫЕ КОДЫ. Владелец уже зафиксировал в src/i18n/ru.ts,
# что IO расшифровывается как MEIO (Multi Echelon Inventory Optimization), а
# INPLAN — это SNP (Supply Network Planning). Завести MEIO и SNP отдельными
# кодами значило бы получить два кода на один модуль — против process-map-b67.
# На слайде 8 презентации MRP модули названы полными именами, поэтому без
# алиасов автоматика молчит: у карты не было бы ни одной внешней системы.
#
# Для карты SNP это безопасно и проверено: в её текстах слова «SNP» и «MEIO»
# как отдельные токены не встречаются ни разу, а сторож — побайтовое совпадение
# пересобранного файла (process-map-3wh.10).
SYSTEM_ALIASES = {"MEIO": "IO", "SNP": "INPLAN"}

# Длинные токены раньше коротких: альтернация в регулярке жадна по порядку.
SYSTEM_TOKENS = tuple(sorted(set(SYSTEM_CODES) | set(SYSTEM_ALIASES), key=len, reverse=True))
SYSTEM_RE = re.compile(r"(?<![A-Za-zА-Яа-я])(" + "|".join(SYSTEM_TOKENS) + r")(?![A-Za-zА-Яа-я])")
DIRECTION_IN_RE = re.compile(r"\bиз\s+модул", re.IGNORECASE)
DIRECTION_OUT_RE = re.compile(r"(\bв\s+модул|Выгрузка\s+в\b)", re.IGNORECASE)


def is_integration(shape: Shape) -> bool:
    """
    Интеграция — серый бокс (заливка A6A6A6) либо текст вида
    «Передача … из/в модуль X». Оба признака взяты из презентации напрямую.
    """
    if shape.fill == INTEGRATION_FILL:
        return True
    return bool(INTEGRATION_TEXT_RE.search(shape.text))


def is_warning(shape: Shape) -> bool:
    """
    Предупреждение — узел, предметом которого является формирование/расчёт
    предупреждения (отдельной заливки для предупреждений в презентации нет,
    единственный однозначный признак — текст).
    """
    return bool(WARNING_TEXT_RE.search(shape.text))


def node_type_for(shape: Shape) -> str:
    if is_integration(shape):
        return "integration"
    if is_warning(shape):
        return "warning"
    return "step"


def detect_system(text: str) -> str | None:
    match = SYSTEM_RE.search(text)
    if match is None:
        return None
    return SYSTEM_ALIASES.get(match.group(1), match.group(1))


def detect_systems(text: str) -> list[str]:
    """
    ВСЕ системы, названные в тексте, в порядке появления и без повторов.

    Нужна там, где одна подпись называет несколько модулей сразу: «Плановые
    заказы из SNP, PS» на слайде 8 — это два входа, а не один. detect_system
    отдаёт только первый и для такого случая не годится.
    """
    seen: set[str] = set()
    codes: list[str] = []
    for match in SYSTEM_RE.finditer(text):
        code = SYSTEM_ALIASES.get(match.group(1), match.group(1))
        if code not in seen:
            seen.add(code)
            codes.append(code)
    return codes


def detect_direction(text: str) -> str | None:
    if DIRECTION_IN_RE.search(text):
        return "in"
    if DIRECTION_OUT_RE.search(text):
        return "out"
    return None


def is_container(shape: Shape) -> bool:
    """Контейнер группы — широкий бесфонный прямоугольник без текста."""
    return (
        shape.kind == "auto"
        and not shape.has_text
        and shape.fill == "noFill"
        and shape.box.width >= CONTAINER_MIN_WIDTH
    )


def is_decor_arrow(shape: Shape) -> bool:
    """
    Мелкая стрелка-иконка между узлами: автофигура без текста, узкая.

    ЗАЛИВКА НЕ ПРОВЕРЯЕТСЯ (process-map-3wh.8). Раньше требовалось fill is None,
    и это работало, пока стрелки были только в презентации SNP. На слайде 8
    презентации MRP такие же стрелки залиты scheme:accent1 — предикат не
    срабатывал, и восемь связей «шаг → шаг» терялись целиком: линиями они там не
    нарисованы вовсе.

    Снятие условия для SNP безопасно и проверено перебором: маленьких
    бестекстовых автофигур на всех шести слайдах ровно шесть, все на слайде 2,
    у всех заливки нет. Сторож — побайтовое совпадение пересобранного файла.
    """
    return (
        shape.kind == "auto"
        and not shape.has_text
        and shape.box.width <= DECOR_ARROW_MAX_WIDTH
    )


def is_caption_long_form(paragraphs: Sequence[str]) -> bool:
    """
    Подпись под узлом — либо перечень выходов, либо развёрнутое описание.
    Описание распознаём по двоеточию, маркерам списка «-» или очень длинному абзацу.
    """
    joined = " ".join(paragraphs)
    if ":" in joined:
        return True
    if any(p.startswith("-") for p in paragraphs):
        return True
    return any(len(p) > 90 for p in paragraphs)


# --------------------------------------------------------------------------------------
# Дедупликация ExternalIO
# --------------------------------------------------------------------------------------

IO_STOP_WORDS = {
    "передача", "передать", "из", "в", "во", "на", "и", "с", "до",
    "модуль", "модуля", "модулей", "выгрузка",
}
IO_MATCH_RATIO = 0.75
IO_MATCH_MIN_COMMON = 2


def io_tokens(text: str) -> set[str]:
    """
    Грубые «стемы» содержательных слов подписи интеграции: нижний регистр, ё→е,
    выброшены стоп-слова, остальное обрезано до 4 символов. Коды систем сохраняются
    целиком. Нужны, чтобы одна и та же передача, описанная на слайде обзора и на
    слайде детализации разными словами, не превратилась в две записи ExternalIO.
    """
    words = re.findall(r"[a-zа-я0-9]+", text.lower().replace("ё", "е"))
    tokens: set[str] = set()
    for word in words:
        if word.upper() in SYSTEM_CODES:
            tokens.add(word)
            continue
        if word in IO_STOP_WORDS or len(word) < 3:
            continue
        tokens.add(word[:4])
    return tokens


def same_transfer(first: str, second: str) -> bool:
    a, b = io_tokens(first), io_tokens(second)
    if not a or not b:
        return first == second
    common = len(a & b)
    if common < IO_MATCH_MIN_COMMON:
        return False
    return common / min(len(a), len(b)) >= IO_MATCH_RATIO


def add_external_io(bucket: list[dict], entry: dict) -> bool:
    """Добавляет ExternalIO, если такая же передача ещё не записана."""
    for existing in bucket:
        if existing["system"] == entry["system"] and same_transfer(existing["label"], entry["label"]):
            return False
    bucket.append(entry)
    return True


# --------------------------------------------------------------------------------------
# Геометрия: группы, подписи, рёбра
# --------------------------------------------------------------------------------------


def find_title_for_container(container: Shape, textboxes: Sequence[Shape]) -> Shape | None:
    """
    Заголовок группы — однострочный TEXT_BOX прямо над контейнером
    (горизонтальное пересечение, верх в окне [container.top - GAP, container.top)).
    Берём самый нижний из подходящих.
    """
    best: Shape | None = None
    for tb in textboxes:
        if tb.consumed_by is not None or len(tb.paragraphs) != 1:
            continue
        if not (container.box.top - GROUP_TITLE_MAX_GAP <= tb.box.top < container.box.top):
            continue
        if tb.box.right <= container.box.left or tb.box.left >= container.box.right:
            continue
        if best is None or tb.box.top > best.box.top:
            best = tb
    return best


def innermost_container(containers: Sequence[tuple[Shape, str]], box: Box) -> str | None:
    """Группа узла — самый маленький контейнер, содержащий центр узла."""
    best: tuple[int, str] | None = None
    for shape, group_id in containers:
        if shape.box.contains_point(box.cx, box.cy):
            if best is None or shape.box.area < best[0]:
                best = (shape.box.area, group_id)
    return best[1] if best else None


def rotate_point(point: tuple[float, float], cx: float, cy: float, degrees: float) -> tuple[float, float]:
    """Поворот точки вокруг центра. Ось Y экранная (вниз), поэтому знак как в DrawingML."""
    if not degrees:
        return point
    angle = math.radians(degrees)
    cos_a, sin_a = math.cos(angle), math.sin(angle)
    dx, dy = point[0] - cx, point[1] - cy
    return (cx + dx * cos_a - dy * sin_a, cy + dx * sin_a + dy * cos_a)


def line_endpoints(shape: Shape) -> tuple[tuple[float, float], tuple[float, float]]:
    """
    Начало и конец линии из left/top/width/height + флагов отражения и ПОВОРОТА.
    Стрелка на конце (a:tailEnd) означает направление начало→конец; стрелка
    на начале (a:headEnd) — обратное.

    ПОВОРОТ (process-map-3wh.18). У повёрнутого коннектора left/top/width/height
    описывают НЕПОВЁРНУТУЮ рамку, а рисуется он повёрнутым вокруг её центра.
    Без доворота концы вычислялись по чужим координатам: у линии [112] слайда 8
    карты MRP геометрия давала начало на «Анализе предупреждений» с расстоянием 0
    — то есть выдуманное ребро, а не потерянную связь. Повёрнутых коннекторов
    хватает и в презентации SNP (11 штук на слайдах 2-5), просто там их
    промахи гасились привязками и вторым проходом.
    """
    x1 = shape.box.right if shape.flip_h else shape.box.left
    x2 = shape.box.left if shape.flip_h else shape.box.right
    y1 = shape.box.bottom if shape.flip_v else shape.box.top
    y2 = shape.box.top if shape.flip_v else shape.box.bottom
    start = (float(x1), float(y1))
    end = (float(x2), float(y2))
    if shape.rot:
        cx, cy = float(shape.box.cx), float(shape.box.cy)
        start = rotate_point(start, cx, cy, shape.rot)
        end = rotate_point(end, cx, cy, shape.rot)
    if shape.head_arrow and not shape.tail_arrow:
        return end, start
    return start, end


def rank_candidates(
    candidates: Sequence[tuple[Box, str]], point: tuple[float, float]
) -> list[tuple[float, str]]:
    best: dict[str, float] = {}
    for box, key in candidates:
        distance = box.distance_to_point(point[0], point[1])
        if key not in best or distance < best[key]:
            best[key] = distance
    # Узел и поглощённая им надпись дают один и тот же key, поэтому конкурентами
    # в проверке однозначности они быть не должны.
    ranked = sorted((distance, key) for key, distance in best.items())
    return ranked


def nearest_target(
    candidates: Sequence[tuple[Box, str]], point: tuple[float, float], snap: float
) -> str | None:
    ranked = rank_candidates(candidates, point)
    if ranked and ranked[0][0] <= snap:
        return ranked[0][1]
    return None


def arrow_edges(
    arrows: Sequence[Shape],
    endpoints: Sequence[tuple[Box, str]],
    snap: float,
) -> list[tuple[str, str]]:
    """
    Декоративные стрелки-иконки как рёбра: левая середина — источник, правая — цель.

    ЗАЧЕМ ОБЩАЯ ФУНКЦИЯ (process-map-3wh.8). Приём был у обзора, но нужен и
    профилю «одиночный слайд»: на слайде 8 презентации MRP связи «шаг → шаг»
    линиями НЕ НАРИСОВАНЫ вовсе — их рисуют восемь таких стрелок. Правило
    проекта «ребро есть, потому что на слайде есть стрелка» соблюдается
    буквально; восстанавливать порядок по абсциссе значило бы догадываться.

    flip_h разворачивает стрелку: с ним источником становится правый край.
    В обеих презентациях отражённых стрелок нет, но правило выражено явно, а не
    оставлено на удачу данных.
    """
    edges: list[tuple[str, str]] = []
    for arrow in arrows:
        left = (float(arrow.box.left), arrow.box.cy)
        right = (float(arrow.box.right), arrow.box.cy)
        start, end = (right, left) if arrow.flip_h else (left, right)
        source = nearest_target(endpoints, start, snap)
        target = nearest_target(endpoints, end, snap)
        if source is None or target is None or source == target:
            # Стрелка между группами внутри одного этапа — не обзорное ребро.
            continue
        edges.append((source, target))
    return edges


def resolve_second_pass(ranked: Sequence[tuple[float, str]], exclude: str) -> str | None:
    """
    Второй проход для линии, у которой разрешён ровно один конец: берём ближайший
    узел в пределах EDGE_SNAP_SECOND, но только если следующий кандидат заметно
    дальше — иначе конец считается неоднозначным и связь не выдумывается.
    """
    filtered = [item for item in ranked if item[1] != exclude]
    if not filtered:
        return None
    distance, key = filtered[0]
    if distance > EDGE_SNAP_SECOND:
        return None
    if len(filtered) > 1 and filtered[1][0] < distance * EDGE_SECOND_RATIO:
        return None
    return key


# --------------------------------------------------------------------------------------
# Обработка слайда детализации
# --------------------------------------------------------------------------------------


def apply_input_enrichment(
    enrichment: dict,
    slide_no: int,
    drafts: list[NodeDraft],
    expanded: set[str],
    ids: IdFactory,
    report: SlideReport,
) -> None:
    """
    Досыпает в колонку входов этапа строки со слайда обзора и проверяет, что
    переформулировки нашли свои узлы (STAGE_INPUT_ENRICHMENT, process-map-qjl).

    Расхождение с презентацией — не предупреждение, а остановка импорта, ровно
    как в apply_owner_decision_edges: если слайд поправят и строка перестанет
    находиться, решение владельца обязано упасть громко, а не рассосаться молча.
    """
    stage_number = enrichment["stage"]
    missed = [short for short, _ in enrichment["expand"] if short not in expanded]
    if missed:
        raise SystemExit(
            f"STAGE_INPUT_ENRICHMENT ({enrichment['task']}): среди входов этапа "
            f"{stage_number} нет строк {', '.join(missed)} — презентация изменилась. "
            f"Импорт остановлен, чтобы решение владельца не потерялось молча: "
            f"обновите список в scripts/import-pptx.py."
        )

    column = [d for d in drafts if d.node_type == "data" and d.direction == "in"]
    if not column:
        raise SystemExit(
            f"STAGE_INPUT_ENRICHMENT ({enrichment['task']}): у этапа {stage_number} нет "
            f"колонки входов — добавлять строки некуда"
        )

    known = {normalize_text(draft.label).casefold() for draft in column}
    # Новые строки встают сразу под колонкой: слайдовая геометрия узла нужна
    # раскладке (layout.ts сидируется slidePosition), а колонка входов — её
    # законное место на слайде детализации.
    anchor = max(column, key=lambda draft: draft.box.bottom)

    for offset, extra in enumerate(enrichment["add"]):
        if normalize_text(extra).casefold() in known:
            raise SystemExit(
                f"STAGE_INPUT_ENRICHMENT ({enrichment['task']}): «{extra}» уже есть среди "
                f"входов этапа {stage_number} — презентация изменилась и строку больше "
                f"добавлять не нужно. Обновите список в scripts/import-pptx.py."
            )
        drafts.append(
            NodeDraft(
                node_id=ids.make(extra, slide_no, enrichment["sid"], offset),
                node_type="data",
                label=extra,
                box=Box(
                    anchor.box.left,
                    anchor.box.bottom + offset * anchor.box.height,
                    anchor.box.width,
                    anchor.box.height,
                ),
                direction="in",
            )
        )
        report.data_nodes += 1
        report.owner_inputs.append(
            f"этап {stage_number}: «{extra}» — ДОБАВЛЕНО со слайда обзора "
            f"({enrichment['task']}), на слайде детализации строки нет"
        )


def build_stage(
    slide_no: int,
    stage_number: int,
    shapes: Sequence[Shape],
    overview_title: str | None,
    ids: IdFactory,
    report: SlideReport,
    seen_signatures: set[tuple],
) -> dict:
    title_shape = next((s for s in shapes if s.kind == "placeholder" and s.has_text), None)
    title = normalize_text(title_shape.text) if title_shape else f"Этап {stage_number}"
    if title_shape:
        title_shape.consumed_by = "заголовок этапа"
    for s in shapes:
        if s.kind == "placeholder" and s.consumed_by is None:
            s.consumed_by = "номер слайда"

    textboxes = [s for s in shapes if s.kind == "textbox" and s.has_text]
    lines = [s for s in shapes if s.kind == "line"]

    # 0. Забытые при копировании слайда фигуры: полностью совпадают с фигурой
    #    предыдущего слайда (тот же shape_id, те же координаты, тот же текст).
    for tb in textboxes:
        signature = (tb.sid, tb.box.left, tb.box.top, tb.box.width, tb.box.height, tb.text)
        if signature in seen_signatures:
            tb.consumed_by = "дубликат фигуры предыдущего слайда"
            report.text_skipped.append(
                f"слайд {slide_no}: [{tb.sid}] «{tb.text[:70]}» — точная копия фигуры "
                f"предыдущего слайда (остаток вёрстки), пропущена"
            )
        else:
            seen_signatures.add(signature)

    # 1. Группы: контейнеры + заголовок над ними.
    containers: list[tuple[Shape, str]] = []
    groups: list[dict] = []
    group_ids: set[str] = set()
    for container in [s for s in shapes if is_container(s)]:
        title_tb = find_title_for_container(container, textboxes)
        if title_tb is None:
            report.questions.append(
                f"слайд {slide_no}: контейнер [{container.sid}] без заголовка — группа не создана"
            )
            continue
        title_tb.consumed_by = "заголовок группы"
        label = title_tb.text
        group_id = slugify(label)
        if group_id in group_ids:
            group_id = f"{group_id}-{container.sid}"
        group_ids.add(group_id)
        containers.append((container, group_id))
        groups.append({"id": group_id, "label": label})
    report.groups = len(groups)

    # 2. Узлы-фигуры: AUTO_SHAPE с текстом (кроме контейнеров и декоративных стрелок).
    node_shapes = [
        s
        for s in shapes
        if s.kind == "auto" and s.has_text and not is_container(s) and not is_decor_arrow(s)
    ]
    drafts: list[NodeDraft] = []
    # Фигура -> узел: по этой карте разрешаются привязки коннекторов stCxn/endCxn.
    node_of_sid: dict[int, str] = {}
    for shape in node_shapes:
        draft = NodeDraft(
            node_id=ids.make(shape.text, slide_no, shape.sid),
            node_type=node_type_for(shape),
            label=shape.text,
            box=shape.box,
            group=innermost_container(containers, shape.box),
        )
        drafts.append(draft)
        node_of_sid[shape.sid] = draft.node_id

    # 2.5. Промоушен текстбокса в узел: в презентации часть вершин графа нарисована
    #      не автофигурой, а надписью. Признак — однострочный текст, в bbox которого
    #      упираются минимум два конца линий (то есть он реально используется как
    #      вершина, а не как подпись).
    endpoint_points: list[tuple[float, float]] = []
    for line in lines:
        start, end = line_endpoints(line)
        endpoint_points.extend((start, end))
    for tb in sorted([t for t in textboxes if t.consumed_by is None], key=Shape.sort_key):
        if len(tb.paragraphs) != 1:
            continue
        hits = sum(
            1 for point in endpoint_points if tb.box.distance_to_point(*point) <= PROMOTE_SNAP
        )
        if hits < PROMOTE_MIN_ENDPOINTS:
            continue
        tb.consumed_by = "узел (промоушен из текстбокса)"
        promoted = NodeDraft(
            node_id=ids.make(tb.text, slide_no, tb.sid),
            node_type=node_type_for(tb),
            label=tb.text,
            box=tb.box,
            group=innermost_container(containers, tb.box),
        )
        drafts.append(promoted)
        node_of_sid[tb.sid] = promoted.node_id
        report.promoted.append(
            f"слайд {slide_no}: [{tb.sid}] «{tb.text[:60]}» — узел (концов линий: {hits})"
        )

    # Текстовые фигуры, поглощённые узлом (подпись-выход или описание), работают
    # как «продолжение» узла: стрелка, упирающаяся в такую надпись, упирается в узел.
    proxies: list[tuple[Box, str]] = []

    # 3. Подписи под узлами → outputs (или description для развёрнутого текста).
    for tb in textboxes:
        if tb.consumed_by is not None:
            continue
        # Текст, дословно повторяющий подпись другого узла этого же слайда, —
        # ссылка на узел, а не новый артефакт: привязка неоднозначна.
        if any(d.label == tb.text for d in drafts):
            report.questions.append(
                f"слайд {slide_no}: [{tb.sid}] «{tb.text[:60]}» дословно повторяет подпись узла "
                f"— как подпись-выход не привязан"
            )
            continue
        best: tuple[float, NodeDraft] | None = None
        for draft in drafts:
            # Тонкие полосы (слайд 5) подписываются сверху, а не снизу.
            if draft.node_type == "data" or draft.box.height < CAPTION_MIN_NODE_HEIGHT:
                continue
            gap = tb.box.top - draft.box.bottom
            if not (CAPTION_MIN_GAP <= gap <= CAPTION_MAX_GAP):
                continue
            overlap = min(tb.box.right, draft.box.right) - max(tb.box.left, draft.box.left)
            if overlap <= 0:
                continue
            covers_node = tb.box.left <= draft.box.left and tb.box.right >= draft.box.right
            coverage = overlap / draft.box.width if covers_node else overlap / tb.box.width
            if coverage < CAPTION_MIN_OVERLAP:
                continue
            score = (gap, abs(tb.box.left - draft.box.left))
            if best is None or score < (best[0], abs(tb.box.left - best[1].box.left)):
                best = (gap, draft)
        if best is None:
            continue
        draft = best[1]
        tb.consumed_by = f"подпись узла {draft.node_id}"
        proxies.append((tb.box, draft.node_id))
        if is_caption_long_form(tb.paragraphs):
            draft.description_parts.extend(tb.paragraphs)
        else:
            draft.outputs.extend(tb.paragraphs)

    # 4. Левая колонка входов → узлы типа data (SPEC §4.2: DataNode входов слева).
    #    Подписи-выходы к этому моменту уже разобраны шагом 3, поэтому здесь
    #    достаточно признака «текстбокс стоит в левом поле слайда».
    #    group у data-узлов не проставляется: входы стоят вне групп (SPEC §4.2).
    #    direction='in' — по происхождению, а не по координате: это и есть
    #    колонка входов слайда (см. NodeDraft.direction, задача process-map-24p).
    #    Часть формулировок берётся со слайда обзора, а не отсюда, — см.
    #    STAGE_INPUT_ENRICHMENT. Замена делается ДО ids.make: иначе id остался бы
    #    слагом старой строки и разошёлся бы с подписью узла.
    enrichment = next((e for e in STAGE_INPUT_ENRICHMENT if e["stage"] == stage_number), None)
    expand = dict(enrichment["expand"]) if enrichment is not None else {}
    expanded: set[str] = set()

    for tb in sorted([t for t in textboxes if t.consumed_by is None], key=Shape.sort_key):
        if tb.box.left > LEFT_MARGIN_LIMIT:
            continue
        tb.consumed_by = "колонка входов"
        unique: list[str] = []
        for para in tb.paragraphs:
            if para not in unique:
                unique.append(para)
        if len(unique) != len(tb.paragraphs):
            report.questions.append(
                f"слайд {slide_no}: в списке входов [{tb.sid}] "
                f"{len(tb.paragraphs) - len(unique)} повторяющихся строк — оставлены уникальные"
            )
        step = tb.box.height / max(len(unique), 1)
        for index, para in enumerate(unique):
            label = expand.get(para, para)
            if label != para:
                expanded.add(para)
                report.owner_inputs.append(
                    f"этап {stage_number}: «{para}» → «{label}» — формулировка взята "
                    f"со слайда обзора ({enrichment['task'] if enrichment else '?'})"
                )
            drafts.append(
                NodeDraft(
                    node_id=ids.make(label, slide_no, tb.sid, index),
                    node_type="data",
                    label=label,
                    box=Box(tb.box.left, int(tb.box.top + index * step), tb.box.width, int(step)),
                    direction="in",
                )
            )
            report.data_nodes += 1

    if enrichment is not None:
        apply_input_enrichment(enrichment, slide_no, drafts, expanded, ids, report)

    # 5. Оставшиеся подписи внутри контейнера → description ближайшего узла группы.
    #    Привязка нестрогая, поэтому каждая попадает в отчёт.
    for tb in sorted([t for t in textboxes if t.consumed_by is None], key=Shape.sort_key):
        group_id = innermost_container(containers, tb.box)
        if group_id is None:
            continue
        siblings = [d for d in drafts if d.group == group_id and d.node_type != "data"]
        if not siblings:
            continue

        def distance(draft: NodeDraft, anchor: Shape = tb) -> tuple[float, str]:
            return (
                math.hypot(draft.box.cx - anchor.box.cx, draft.box.cy - anchor.box.cy),
                draft.node_id,
            )

        # В презентации такие подписи стоят НАД своим узлом (ряды «повод — действие»
        # на слайде 5), поэтому сначала ищем ближайший узел ниже подписи.
        below = [d for d in siblings if d.box.cy > tb.box.cy]
        best_draft = min(below or siblings, key=distance)
        dist = distance(best_draft)[0]
        if len(siblings) > 1 and dist > LOOSE_DESC_MAX_DIST:
            continue
        tb.consumed_by = f"описание узла {best_draft.node_id}"
        proxies.append((tb.box, best_draft.node_id))
        best_draft.description_parts.extend(tb.paragraphs)
        report.loose_attachments.append(
            f"слайд {slide_no}: [{tb.sid}] «{tb.text[:70]}» → description узла «{best_draft.node_id}»"
        )

    # 6. Рёбра: сначала явные привязки коннектора, затем геометрия (два прохода).
    #
    #    ПОЧЕМУ ПРИВЯЗКИ ПЕРВЫМИ (process-map-3wh.16). Считалось, что в этой
    #    презентации привязок нет и связи приходится выводить геометрически. Это
    #    неверно: из 66 линий 47 имеют ОБЕ привязки и лишь 10 не имеют ни одной —
    #    импортёр их просто не читал. Привязка точнее геометрии: она указывает на
    #    фигуру, а не на точку, и не врёт у растянутых коннекторов.
    #
    #    Привязка к КОНТЕЙНЕРУ группы разрешается в ближайший шаг этой группы:
    #    ребро в модели соединяет узлы, а контейнер узлом не является. Веерную
    #    связь «в группу целиком» один коннектор выразить не может — для неё
    #    остаётся OWNER_DECISION_EDGES (проверено: линия [146] слайда 5 даёт одно
    #    ребро из четырёх, объявленных владельцем).
    by_id = {d.node_id: d for d in drafts}
    node_boxes = [(d.box, d.node_id) for d in drafts if d.node_type != "data"]
    node_boxes.extend(proxies)
    members_of = {
        container.sid: [
            (d.box, d.node_id)
            for d in drafts
            if d.group == group_id and d.node_type != "data"
        ]
        for container, group_id in containers
    }

    def by_binding(sid: int | None, point: tuple[float, float]) -> str | None:
        """Привязка коннектора -> узел; привязка к контейнеру -> ближайший шаг в нём."""
        if sid is None:
            return None
        if sid in node_of_sid:
            return node_of_sid[sid]
        return nearest_target(members_of.get(sid, []), point, EDGE_SNAP_DETAIL)

    edges: list[dict] = []
    seen_pairs: set[tuple[str, str]] = set()
    report.lines_total = len(lines)
    for line in lines:
        start, end = line_endpoints(line)
        source = by_binding(line.start_sid, start)
        target = by_binding(line.end_sid, end)
        bound = source is not None or target is not None
        ranked_start = rank_candidates(node_boxes, start)
        ranked_end = rank_candidates(node_boxes, end)
        if source is None and ranked_start and ranked_start[0][0] <= EDGE_SNAP_DETAIL:
            source = ranked_start[0][1]
        if target is None and ranked_end and ranked_end[0][0] <= EDGE_SNAP_DETAIL:
            target = ranked_end[0][1]
        if source is not None and target is None:
            target = resolve_second_pass(ranked_end, source)
        elif target is not None and source is None:
            source = resolve_second_pass(ranked_start, target)
        if source is None or target is None or source == target:
            reason = "петля" if source is not None and source == target else "конец не определён"
            report.lines_skipped.append(f"слайд {slide_no}: линия [{line.sid}] — {reason}")
            continue
        if (source, target) in seen_pairs:
            continue
        seen_pairs.add((source, target))
        if bound:
            report.cxn_edges.append(f"слайд {slide_no}: [{line.sid}] {source} → {target}")
        endpoint_types = {by_id[source].node_type, by_id[target].node_type}
        edges.append(
            {
                "id": f"e-{source}--{target}",
                "source": source,
                "target": target,
                "kind": "integration" if "integration" in endpoint_types else "process",
            }
        )
    report.edges = len(edges)

    # 7. Внешние системы этапа (ExternalIO): из текста интеграций и узлов публикации.
    inputs: list[dict] = []
    outputs: list[dict] = []
    for draft in drafts:
        if draft.node_type == "integration" and detect_system(draft.label) is None:
            report.questions.append(
                f"слайд {slide_no}: интеграция «{draft.label[:60]}» — система в тексте не названа"
            )
        for text in [draft.label] + draft.outputs:
            system = detect_system(text)
            direction = detect_direction(text)
            if system is None or direction is None:
                continue
            entry = {"system": system, "label": text, "stage": stage_number, "direction": direction}
            bucket = inputs if direction == "in" else outputs
            if add_external_io(bucket, entry) and draft.system is None:
                draft.system = system

    # 7b. Узел с проставленным кодом системы — интеграция, даже если заливка
    #     обычная (решение владельца, process-map-7v1).
    #
    #     ПОЧЕМУ ПРАВИЛО, А НЕ СПИСОК. Признак «серая заливка A6A6A6» есть ровно
    #     у четырёх фигур на всю презентацию, и все четыре — входящие интеграции
    #     слайдов 3-4. Исходящие интеграции слайда 5 нарисованы штатной заливкой
    #     шага, а слово «модуль» из их текста убрано («Передача ограниченного
    #     прогноза в DP»), хотя на слайде обзора тот же шаг записан «в модуль DP».
    #     То есть оба признака is_integration промахиваются по вине вёрстки.
    #
    #     Код системы при этом проставляется от подписи-выгрузки и служит честным
    #     признаком: узел назвал внешнюю систему и направление — значит он на
    #     границе с ней. На текущих данных правило даёт ровно те четыре узла,
    #     которые назвал владелец, и ни одного лишнего: код системы несут восемь
    #     узлов, четыре из них уже были интеграциями.
    for draft in drafts:
        if draft.system is not None and draft.node_type == "step":
            draft.node_type = "integration"
            report.promoted_integrations.append(
                f"слайд {slide_no}: «{draft.label[:60]}» — step → integration "
                f"(код системы {draft.system}, заливка обычная)"
            )

    # 8. Отчёт по потерянным текстовым фигурам.
    for tb in textboxes:
        if tb.consumed_by is None:
            report.text_skipped.append(f"слайд {slide_no}: [{tb.sid}] «{tb.text[:80]}»")

    report.nodes = len(drafts)

    short_title = re.split(r"\s/|/\s|\s\+\s", title)[0].strip()
    return {
        "id": slugify(f"stage-{stage_number}-{overview_title or short_title}"),
        "number": stage_number,
        "title": title,
        "shortTitle": short_title,
        "keyOutputs": [],
        "warningsCount": 0,
        "groups": groups,
        "nodes": [serialize_node(d) for d in sorted(drafts, key=lambda d: (d.box.top, d.box.left, d.node_id))],
        "edges": edges,
        "inputs": inputs,
        "outputs": outputs,
    }


def serialize_node(draft: NodeDraft) -> dict:
    node: dict = {"id": draft.node_id, "type": draft.node_type, "label": draft.label}
    if draft.description:
        node["description"] = draft.description
    if draft.group:
        node["group"] = draft.group
    if draft.direction:
        node["direction"] = draft.direction
    if draft.inputs:
        node["inputs"] = draft.inputs
    if draft.outputs:
        node["outputs"] = draft.outputs
    if draft.system:
        node["system"] = draft.system
    slide_position = {
        "x": round(draft.box.left / EMU_PER_PX),
        "y": round(draft.box.top / EMU_PER_PX),
    }
    # position — то, что покажет приложение; его перезапишет `npm run layout`.
    # slidePosition — та же геометрия слайда, но НАВСЕГДА: раскладка сидируется
    # ею, а не собственным прошлым результатом (SPEC §3, задача process-map-cxn).
    # Два разных словаря, а не один и тот же объект: иначе правка одного поля
    # молча меняла бы второе.
    node["position"] = dict(slide_position)
    node["slidePosition"] = slide_position
    return node


# --------------------------------------------------------------------------------------
# Обработка слайда обзора
# --------------------------------------------------------------------------------------


@dataclass
class OverviewData:
    titles: list[str]
    output_blocks: list[list[tuple[list[str], Box, int]]]
    group_labels: list[list[str]]
    systems: list[dict]
    edges: list[dict]
    report: SlideReport


def build_overview(shapes: Sequence[Shape], report: SlideReport) -> OverviewData:
    textboxes = [s for s in shapes if s.kind == "textbox" and s.has_text]
    for s in shapes:
        if s.kind == "placeholder":
            s.consumed_by = "служебный текст слайда"

    containers = [s for s in shapes if is_container(s)]
    # Порядок этапов: верхний ряд, затем нижний слева направо.
    containers.sort(key=lambda s: (round(s.box.top / 1_000_000), s.box.left))

    titles: list[str] = []
    output_blocks: list[list[tuple[list[str], Box, int]]] = []
    group_labels: list[list[str]] = []

    boxed = [s for s in shapes if s.kind == "auto" and s.has_text]
    group_boxes = [s for s in boxed if s.fill == "scheme:accent1"]
    system_boxes = [s for s in boxed if s.fill in ("scheme:accent2", "scheme:bg1")]

    for container in containers:
        title_tb = find_title_for_container(container, textboxes)
        if title_tb is not None:
            title_tb.consumed_by = "заголовок этапа обзора"
            titles.append(title_tb.text)
        else:
            titles.append("")
            report.questions.append(
                f"слайд 2: контейнер этапа [{container.sid}] без подписи-заголовка "
                f"— заголовок берётся со слайда детализации"
            )

        group_labels.append(
            [
                b.text
                for b in sorted(group_boxes, key=Shape.sort_key)
                if container.box.contains_point(b.box.cx, b.box.cy)
            ]
        )

        lo = container.box.top + KEY_OUTPUT_TOP_OFFSET
        hi = container.box.bottom + KEY_OUTPUT_BOTTOM_OFFSET
        blocks: list[tuple[list[str], Box, int]] = []
        for tb in sorted(
            [t for t in textboxes if t.consumed_by is None], key=lambda s: (s.box.left, s.box.top)
        ):
            if not (lo <= tb.box.top <= hi):
                continue
            if tb.box.right <= container.box.left or tb.box.left >= container.box.right:
                continue
            tb.consumed_by = "блок выходов этапа"
            blocks.append((list(tb.paragraphs), tb.box, tb.sid))
        output_blocks.append(blocks)

    for b in group_boxes:
        if b.consumed_by is None:
            b.consumed_by = "бокс группы обзора"

    systems: list[dict] = []
    for shape in sorted(system_boxes, key=Shape.sort_key):
        code = detect_system(shape.text)
        if code is None:
            report.questions.append(
                f"слайд 2: бокс [{shape.sid}] «{shape.text[:60]}» — система не распознана"
            )
            continue
        systems.append(
            {
                "code": code,
                "label": shape.text,
                "direction": detect_direction(shape.text),
                "box": shape.box,
                "sid": shape.sid,
            }
        )
        shape.consumed_by = "внешняя система обзора"

    endpoints: list[tuple[Box, str]] = []
    for index, container in enumerate(containers, start=1):
        endpoints.append((container.box, f"@stage{index}"))
    for system in systems:
        endpoints.append((system["box"], f"@sys:{system['code']}:{system['sid']}"))

    connectors = [s for s in shapes if s.kind == "line"]
    arrows = [s for s in shapes if is_decor_arrow(s)]
    report.lines_total = len(connectors) + len(arrows)

    raw_edges: list[tuple[str, str]] = []
    for line in connectors:
        start, end = line_endpoints(line)
        source = nearest_target(endpoints, start, EDGE_SNAP_OVERVIEW)
        target = nearest_target(endpoints, end, EDGE_SNAP_OVERVIEW)
        if source is None or target is None or source == target:
            report.lines_skipped.append(
                f"слайд 2: линия [{line.sid}] — концы не привязались к этапу/системе"
            )
            continue
        raw_edges.append((source, target))

    raw_edges.extend(arrow_edges(arrows, endpoints, EDGE_SNAP_OVERVIEW))

    for tb in textboxes:
        if tb.consumed_by is None:
            report.text_skipped.append(f"слайд 2: [{tb.sid}] «{tb.text[:80]}»")

    return OverviewData(
        titles=titles,
        output_blocks=output_blocks,
        group_labels=group_labels,
        systems=systems,
        edges=[{"source": s, "target": t} for s, t in raw_edges],
        report=report,
    )


def block_items_start(paragraphs: Sequence[str]) -> int:
    """
    С какого абзаца в блоке выходов начинаются сами пункты.

    Блок может открываться заголовком — на слайде 2 это «Опубликованные планы:».
    Признак один: первый абзац оканчивается двоеточием. Разметка его
    подтверждает (единственный во всей презентации абзац без буллита и с нулевым
    отступом на фоне буллитованных соседей), но опираться на неё нельзя —
    двоеточие надёжнее и не зависит от того, как редактор записал отступы.

    Функция одна на ОБА применения — выбор keyOutputs и создание узлов-выходов.
    Раньше правило жило только в первом, и заголовок уезжал в данные отдельной
    карточкой «Опубликованные планы:» — единственным узлом документа, чья
    подпись оканчивалась двоеточием (process-map-t9j).
    """
    return 1 if paragraphs and paragraphs[0].endswith(":") else 0


def choose_key_outputs(
    blocks: Sequence[tuple[list[str], Box, int]],
    inbound_labels: AbstractSet[str],
) -> list[str]:
    """
    Ключевые выходы этапа. Если среди блоков есть перечень с заголовком-двоеточием
    («Опубликованные планы:»), берём его пункты — это и есть выходы этапа.
    Иначе — первые MAX_KEY_OUTPUTS абзацев в порядке чтения.

    ЧТО ОТБРАСЫВАЕТСЯ (process-map-tze). Пункт, который на слайде детализации
    этого же этапа нарисован ВХОДОМ, ключевым выходом не считается: иначе
    карточка обзора называет выходом то, что на экране этапа стоит в левой
    колонке, и карта противоречит сама себе на двух соседних экранах.

    Так вышло из презентации, и это не её ошибка: блоки выходов слайда 2 стоят
    под боксами ГРУПП, а не под этапом (process-map-2of). «Прогноз продаж в
    требуемой гранулярности» — выход группы «Преобразование прогноза к
    требуемой гранулярности» и одновременно вход следующей группы. Для этапа
    целиком это вход, для группы внутри него — выход.

    Отбор идёт ДО среза по MAX_KEY_OUTPUTS: иначе отброшенный пункт занимал бы
    место, и на карточке оказалось бы меньше выходов, чем мест.
    """

    def keep(paragraphs: Iterable[str]) -> list[str]:
        return [p for p in paragraphs if p.casefold() not in inbound_labels]

    for paragraphs, _box, _sid in blocks:
        start = block_items_start(paragraphs)
        if start:
            return keep(paragraphs[start:])[:MAX_KEY_OUTPUTS]
    flat = [p for paragraphs, _box, _sid in blocks for p in paragraphs]
    return keep(flat)[:MAX_KEY_OUTPUTS]


# --------------------------------------------------------------------------------------
# Сборка документа
# --------------------------------------------------------------------------------------


def count_warnings(stage: dict) -> int:
    """
    Число предупреждений этапа — по содержанию слайда (перечень типов
    предупреждений), а не по числу узлов типа warning.
    """
    distinct: set[str] = set()
    for node in stage["nodes"]:
        candidates = [node["label"], *node.get("outputs", [])]
        candidates.extend((node.get("description") or "").split("\n"))
        for text in candidates:
            text = text.strip()
            if text and WARNING_ITEM_RE.match(text):
                distinct.add(text)
    node_warnings = sum(1 for node in stage["nodes"] if node["type"] == "warning")
    return max(len(distinct), node_warnings)


def is_artifact_box(shape: Shape) -> bool:
    """
    Плашка-артефакт: вход или выход процесса, нарисованный автофигурой.

    ЗАЧЕМ ОТДЕЛЬНЫЙ ПРЕДИКАТ. На слайдах детализации SNP колонка входов — это
    надписи (textbox), и разбор ищет их по признаку «текстбокс в левом поле».
    На слайде 8 презентации MRP те же входы и выход нарисованы АВТОФИГУРАМИ с
    заливкой scheme:accent2. Без этого предиката они стали бы обычными шагами:
    прогон существующего build_stage давал 17 узлов и НОЛЬ data-узлов.
    """
    return shape.kind == "auto" and shape.has_text and shape.fill == ARTIFACT_FILL


def build_single_slide_map(
    collisions: dict[str, int] | None,
    spec: MapSpec,
) -> tuple[dict, list[SlideReport], list[str], Counter[str]]:
    """
    Профиль «одиночный слайд»: вся карта собирается с ОДНОГО слайда.

    ЧЕМ ОТЛИЧАЕТСЯ ОТ build_process_map. Там два уровня слайдов: обзор со
    стадиями и четыре слайда детализации. Здесь уровня «обзор» нет вовсе, и
    четыре КОНТЕЙНЕРА слайда становятся четырьмя этапами карты (решение
    владельца). Группы внутри этапа не заводятся: группа СТАЛА этапом, и
    рисовать внутри него рамку с тем же заголовком значило бы дублировать.

    Общего с build_stage — почти вся геометрия: is_container, поиск заголовка,
    innermost_container, классификация узлов, разбор подписей, двухпроходное
    сопоставление концов линий. Своё — только «где брать этапы» и «где брать
    входы и выходы».
    """
    if not spec.pptx.exists():
        raise SystemExit(f"Не найдена презентация: {spec.pptx}")
    presentation = Presentation(str(spec.pptx))
    if int(presentation.slide_width) != SLIDE_WIDTH_EMU:
        raise SystemExit(
            f"Неожиданная ширина слайда {presentation.slide_width} EMU (ожидалось {SLIDE_WIDTH_EMU})"
        )
    slides = list(presentation.slides)
    if len(slides) != spec.slides:
        raise SystemExit(f"Ожидалось {spec.slides} слайдов, найдено {len(slides)}")
    if spec.slide_index is None:
        raise SystemExit(f"У карты «{spec.key}» профиль single-slide, но slide_index не задан")

    slide_no = spec.slide_index + 1
    shapes = read_slide(slides[spec.slide_index])
    report = SlideReport(slide_no=slide_no)
    questions: list[str] = []

    textboxes = [s for s in shapes if s.kind == "textbox" and s.has_text]
    lines = [s for s in shapes if s.kind == "line"]
    arrows = [s for s in shapes if is_decor_arrow(s)]
    for s in shapes:
        if s.kind == "placeholder":
            s.consumed_by = "заголовок слайда"

    # 1. Контейнеры -> этапы. Порядок «верхний ряд, затем нижний слева направо» —
    #    тот же, которым build_overview упорядочивает стадии обзора.
    containers = sorted(
        [s for s in shapes if is_container(s)],
        key=lambda s: (round(s.box.top / 1_000_000), s.box.left),
    )
    if len(containers) != STAGE_COUNT:
        raise SystemExit(
            f"слайд {slide_no}: ожидалось {STAGE_COUNT} контейнеров-этапов, "
            f"найдено {len(containers)}"
        )

    stage_containers: list[tuple[Shape, str]] = []
    stage_meta: list[dict] = []
    for index, container in enumerate(containers):
        title_tb = find_title_for_container(container, textboxes)
        if title_tb is None:
            # Не вопрос в отчёт, а остановка: в этом профиле контейнер ЭТО этап,
            # и безымянный этап собирать нельзя.
            raise SystemExit(
                f"слайд {slide_no}: контейнер [{container.sid}] без заголовка — "
                f"этап собрать нельзя"
            )
        title_tb.consumed_by = "заголовок этапа"
        title = title_tb.text
        short_title = re.split(r"\s/|/\s|\s\+\s", title)[0].strip()
        number = index + 1
        stage_id = slugify(f"stage-{number}-{short_title}")
        stage_containers.append((container, stage_id))
        stage_meta.append(
            {"id": stage_id, "number": number, "title": title, "shortTitle": short_title}
        )
    report.groups = len(stage_meta)

    # 2. Узлы-шаги. Плашки-артефакты исключены: они станут data-узлами шагом 3.
    ids = IdFactory(collisions)
    drafts: list[NodeDraft] = []
    stage_of_node: dict[str, str] = {}
    sid_of_node: dict[int, str] = {}
    node_shapes = [
        s
        for s in shapes
        if s.kind == "auto"
        and s.has_text
        and not is_container(s)
        and not is_decor_arrow(s)
        and not is_artifact_box(s)
    ]
    for shape in node_shapes:
        stage_id = innermost_container(stage_containers, shape.box)
        if stage_id is None:
            raise SystemExit(
                f"слайд {slide_no}: узел «{shape.text[:60]}» не попал ни в один контейнер-этап"
            )
        draft = NodeDraft(
            node_id=ids.make(shape.text, slide_no, shape.sid),
            node_type=node_type_for(shape),
            label=shape.text,
            box=shape.box,
        )
        drafts.append(draft)
        stage_of_node[draft.node_id] = stage_id
        sid_of_node[shape.sid] = draft.node_id
    report.nodes = len(drafts)

    # 3. Плашки-артефакты -> data-узлы.
    #    НАПРАВЛЕНИЕ ПО ПРОИСХОЖДЕНИЮ, а не по координате (process-map-24p):
    #    левое поле слайда — это колонка входов. Для плашки вне левого поля
    #    правило «значит выход» было бы геометрией, поэтому шагом 6 оно
    #    подкрепляется вторым независимым признаком: такая плашка обязана быть
    #    целью хотя бы одной связи.
    artifacts: list[tuple[Shape, NodeDraft]] = []
    for shape in [s for s in shapes if is_artifact_box(s)]:
        direction = "in" if shape.box.left <= LEFT_MARGIN_LIMIT else "out"
        draft = NodeDraft(
            node_id=ids.make(shape.text, slide_no, shape.sid),
            node_type="data",
            label=shape.text,
            box=shape.box,
            direction=direction,
        )
        drafts.append(draft)
        artifacts.append((shape, draft))
        sid_of_node[shape.sid] = draft.node_id
        report.data_nodes += 1
    if not artifacts:
        raise SystemExit(
            f"слайд {slide_no}: не найдено ни одной плашки-артефакта (заливка {ARTIFACT_FILL})"
        )

    # 4. Подписи под шагами -> inputs узла (решение владельца, process-map-3wh.1).
    #    На слайде 8 это перечни ИСХОДНЫХ ДАННЫХ («Спецификации (BOM) с уровнями
    #    и нормами расхода», «Остатки и страховые запасы»). Правило SNP отправило
    #    бы два из трёх в outputs, где перечень входов был бы враньём.
    proxies: list[tuple[Box, str]] = []
    step_drafts = [d for d in drafts if d.node_type != "data"]
    for tb in sorted([t for t in textboxes if t.consumed_by is None], key=Shape.sort_key):
        best: tuple[float, NodeDraft] | None = None
        for draft in step_drafts:
            if draft.box.height < CAPTION_MIN_NODE_HEIGHT:
                continue
            gap = tb.box.top - draft.box.bottom
            if not (CAPTION_MIN_GAP <= gap <= CAPTION_MAX_GAP):
                continue
            overlap = min(tb.box.right, draft.box.right) - max(tb.box.left, draft.box.left)
            if overlap <= 0:
                continue
            covers_node = tb.box.left <= draft.box.left and tb.box.right >= draft.box.right
            coverage = overlap / draft.box.width if covers_node else overlap / tb.box.width
            if coverage < CAPTION_MIN_OVERLAP:
                continue
            score = (gap, abs(tb.box.left - draft.box.left))
            if best is None or score < (best[0], abs(tb.box.left - best[1].box.left)):
                best = (gap, draft)
        if best is None:
            continue
        draft = best[1]
        tb.consumed_by = f"входы узла {draft.node_id}"
        proxies.append((tb.box, draft.node_id))
        # Заголовок перечня («Параметры закупки:») — не исходные данные, а
        # подпись над ними. Та же функция и тот же дефект, что в process-map-t9j:
        # там заголовок уезжал в данные отдельной карточкой, здесь — отдельным
        # входом. Срез не применяется, если после него ничего не осталось:
        # подпись из одного заголовка лучше сохранить целиком, чем потерять.
        items = tb.paragraphs[block_items_start(tb.paragraphs) :]
        draft.inputs.extend(items or tb.paragraphs)

    # 5. Рёбра — три источника в жёстком порядке.
    node_boxes: list[tuple[Box, str]] = [(d.box, d.node_id) for d in drafts]
    node_boxes.extend(proxies)
    step_boxes = [(d.box, d.node_id) for d in step_drafts]
    members_of: dict[int, list[tuple[Box, str]]] = {}
    for container, stage_id in stage_containers:
        members_of[container.sid] = [
            (d.box, d.node_id) for d in step_drafts if stage_of_node.get(d.node_id) == stage_id
        ]

    raw_edges: list[tuple[str, str]] = []
    seen_pairs: set[tuple[str, str]] = set()

    def remember(source: str | None, target: str | None, sid: int) -> None:
        if source is None or target is None or source == target:
            reason = "петля" if source is not None and source == target else "конец не определён"
            report.lines_skipped.append(f"слайд {slide_no}: линия [{sid}] — {reason}")
            return
        if (source, target) in seen_pairs:
            return
        seen_pairs.add((source, target))
        raw_edges.append((source, target))

    # 5a. Декоративные стрелки -> связи «шаг -> шаг». Линиями они на слайде НЕ
    #     нарисованы вовсе: порядок шагов внутри этапа держится только на них.
    #     Восстанавливать его по абсциссе значило бы догадываться.
    for source, target in arrow_edges(arrows, step_boxes, EDGE_SNAP_DETAIL):
        if (source, target) not in seen_pairs:
            seen_pairs.add((source, target))
            raw_edges.append((source, target))
    report.lines_total = len(lines) + len(arrows)

    def by_binding(sid: int | None, point: tuple[float, float]) -> str | None:
        """Привязка коннектора -> узел. Привязка к КОНТЕЙНЕРУ — ближайший шаг в нём."""
        if sid is None:
            return None
        if sid in sid_of_node:
            return sid_of_node[sid]
        return nearest_target(members_of.get(sid, []), point, EDGE_SNAP_DETAIL)

    for line in lines:
        start, end = line_endpoints(line)
        # 5b. Привязки читаются ПЕРВЫМИ: у линии обратной связи слайда 8 bbox
        #     уходит ниже кромки слайда (растянутый коннектор), и геометрия её
        #     теряет, а привязка — нет.
        source = by_binding(line.start_sid, start)
        target = by_binding(line.end_sid, end)
        # 5c. Чего привязки не дали — добирается геометрией, двумя проходами.
        ranked_start = rank_candidates(node_boxes, start)
        ranked_end = rank_candidates(node_boxes, end)
        if source is None and ranked_start and ranked_start[0][0] <= EDGE_SNAP_DETAIL:
            source = ranked_start[0][1]
        if target is None and ranked_end and ranked_end[0][0] <= EDGE_SNAP_DETAIL:
            target = ranked_end[0][1]
        if source is not None and target is None:
            target = resolve_second_pass(ranked_end, source)
        elif target is not None and source is None:
            source = resolve_second_pass(ranked_start, target)
        remember(source, target, line.sid)

    # 6. Этап data-узла — этап шага, с которым он связан.
    by_id = {d.node_id: d for d in drafts}
    for shape, draft in artifacts:
        partners = [t for s, t in raw_edges if s == draft.node_id]
        partners += [s for s, t in raw_edges if t == draft.node_id]
        steps = [p for p in partners if by_id[p].node_type != "data"]
        if not steps:
            raise SystemExit(
                f"слайд {slide_no}: плашка-артефакт [{shape.sid}] «{draft.label[:50]}» "
                f"не связана ни с одним шагом — направление «{draft.direction}» недоказуемо"
            )
        if draft.direction == "out" and not any(t == draft.node_id for _, t in raw_edges):
            raise SystemExit(
                f"слайд {slide_no}: плашка [{shape.sid}] «{draft.label[:50]}» признана выходом "
                f"по положению, но не является целью ни одной связи"
            )
        stage_of_node[draft.node_id] = stage_of_node[steps[0]]

    # 6b. Внешние системы этапа. Направление берётся ИЗ ПРОИСХОЖДЕНИЯ ФИГУРЫ
    #     (draft.direction), а не из слова в тексте: detect_direction ищет
    #     «из модул…»/«в модул…», а на слайде 8 написано «…из SNP, PS» — то есть
    #     без правила «плашка левой колонки это вход» у карты не было бы ни одной
    #     внешней системы, ни одного свимлейна и ни одной записи в stage.inputs.
    #
    #     Система выходной плашки в тексте не названа («…в систему исполнения
    #     закупок»), и угадывать её запрещено (CLAUDE.md). Такая плашка остаётся
    #     обычным data-узлом без ExternalIO — это отказ выдумывать, а не заглушка.
    stage_io: dict[str, tuple[list[dict], list[dict]]] = {
        meta["id"]: ([], []) for meta in stage_meta
    }
    number_of = {meta["id"]: meta["number"] for meta in stage_meta}
    for shape, draft in artifacts:
        stage_id = stage_of_node[draft.node_id]
        codes = detect_systems(draft.label)
        if not codes:
            questions.append(
                f"слайд {slide_no}: артефакт «{draft.label[:60]}» — система в тексте не названа, "
                f"внешняя система для него не заведена"
            )
            continue
        inputs_bucket, outputs_bucket = stage_io[stage_id]
        bucket = inputs_bucket if draft.direction == "in" else outputs_bucket
        for code in codes:
            add_external_io(
                bucket,
                {
                    "system": code,
                    "label": draft.label,
                    "stage": number_of[stage_id],
                    "direction": draft.direction,
                },
            )
        if draft.system is None:
            draft.system = codes[0]

    # 7. Рёбра внутри этапа -> stage.edges, между этапами -> overviewEdges.
    stage_edges: dict[str, list[dict]] = {meta["id"]: [] for meta in stage_meta}
    overview_edges: list[dict] = []
    seen_overview: set[tuple[str, str]] = set()
    for source, target in raw_edges:
        source_stage = stage_of_node[source]
        target_stage = stage_of_node[target]
        kind = (
            "integration"
            if "integration" in {by_id[source].node_type, by_id[target].node_type}
            else "process"
        )
        if source_stage == target_stage:
            stage_edges[source_stage].append(
                {"id": f"e-{source}--{target}", "source": source, "target": target, "kind": kind}
            )
            continue
        pair = (source_stage, target_stage)
        if pair in seen_overview:
            continue
        seen_overview.add(pair)
        overview_edges.append(
            {
                "id": f"ov-{source_stage}--{target_stage}",
                "source": source_stage,
                "target": target_stage,
                "kind": "process",
            }
        )
    # 7b. Связи «система -> этап» на обзоре. Выводятся прямо из ExternalIO, а не
    #     из линий: на слайде нет уровня обзора, но принадлежность известна точно
    #     — артефакт лежит в своём этапе и называет свою систему. Формат концов
    #     тот же, что у карты SNP: код системы как идентификатор.
    for meta in stage_meta:
        stage_id = meta["id"]
        inputs_bucket, outputs_bucket = stage_io[stage_id]
        for entry in inputs_bucket:
            overview_edges.append(
                {
                    "id": f"ov-{entry['system']}--{stage_id}",
                    "source": entry["system"],
                    "target": stage_id,
                    "kind": "integration",
                }
            )
        for entry in outputs_bucket:
            overview_edges.append(
                {
                    "id": f"ov-{stage_id}--{entry['system']}",
                    "source": stage_id,
                    "target": entry["system"],
                    "kind": "integration",
                }
            )

    report.edges = sum(len(items) for items in stage_edges.values())

    # 8. Сборка этапов.
    stages: list[dict] = []
    for meta in stage_meta:
        stage_id = meta["id"]
        members = [d for d in drafts if stage_of_node.get(d.node_id) == stage_id]
        key_outputs = [d.label for d in members if d.direction == "out"][:MAX_KEY_OUTPUTS]
        stages.append(
            {
                "id": stage_id,
                "number": meta["number"],
                "title": meta["title"],
                "shortTitle": meta["shortTitle"],
                "keyOutputs": key_outputs,
                # warningsCount НЕ ПРОСТАВЛЯЕТСЯ. У SNP счётчик означает «сколько
                # типов предупреждений формирует этап»; переносить эту семантику
                # на шаг «Анализ предупреждений», который планировщик выполняет
                # сам, значило бы соврать. Поле необязательное (SPEC §3).
                "groups": [],
                "nodes": [
                    serialize_node(d)
                    for d in sorted(members, key=lambda d: (d.box.top, d.box.left, d.node_id))
                ],
                "edges": stage_edges[stage_id],
                "inputs": stage_io[stage_id][0],
                "outputs": stage_io[stage_id][1],
            }
        )

    for tb in textboxes:
        if tb.consumed_by is None:
            report.text_skipped.append(f"слайд {slide_no}: [{tb.sid}] «{tb.text[:80]}»")

    process_map = {
        "version": MAP_VERSION,
        "id": spec.map_id,
        "updatedAt": spec.updated_at,
        "title": spec.title,
        "moduleLabel": spec.module_label,
        "stages": stages,
        "overviewEdges": overview_edges,
    }
    return process_map, [report], questions, ids.counts


def build_process_map(
    collisions: dict[str, int] | None,
    spec: MapSpec,
) -> tuple[dict, list[SlideReport], list[str], Counter[str]]:
    if not spec.pptx.exists():
        raise SystemExit(f"Не найдена презентация: {spec.pptx}")

    presentation = Presentation(str(spec.pptx))
    if int(presentation.slide_width) != SLIDE_WIDTH_EMU:
        raise SystemExit(
            f"Неожиданная ширина слайда {presentation.slide_width} EMU (ожидалось {SLIDE_WIDTH_EMU})"
        )
    slides = list(presentation.slides)
    if len(slides) != spec.slides:
        raise SystemExit(f"Ожидалось {spec.slides} слайдов, найдено {len(slides)}")

    reports: list[SlideReport] = []
    questions: list[str] = []

    overview_report = SlideReport(slide_no=2)
    overview_shapes = read_slide(slides[1])
    overview = build_overview(overview_shapes, overview_report)

    ids = IdFactory(collisions)
    seen_signatures: set[tuple] = set()
    stages: list[dict] = []
    for index in range(STAGE_COUNT):
        slide_no = index + 3
        report = SlideReport(slide_no=slide_no)
        stage = build_stage(
            slide_no=slide_no,
            stage_number=index + 1,
            shapes=read_slide(slides[index + 2]),
            overview_title=overview.titles[index] if index < len(overview.titles) else None,
            ids=ids,
            report=report,
            seen_signatures=seen_signatures,
        )
        stages.append(stage)
        reports.append(report)

    # Рёбра, которых в презентации нет (OWNER_DECISION_EDGES) — досыпаются сразу
    # после разбора слайдов детализации, до сверок и отчётов: изолированные узлы
    # и целостность считаются уже по итоговому набору рёбер. reports здесь ещё
    # содержит ровно 4 отчёта этапов, отчёт обзора вставляется в начало ниже.
    #
    # Только во ВТОРОЙ фазе: в первой (collisions=None) IdFactory выдаёт
    # временные id вида «base~N» ради подсчёта коллизий, и сверять с ними id,
    # названные в решении владельца, бессмысленно — не нашлось бы ни одного.
    # Результат первой фазы всё равно выбрасывается, кроме карты коллизий, а
    # рёбра решения новых id не создают, так что пропуск ни на что не влияет.
    if collisions is not None:
        apply_owner_decision_edges(stages, reports)
        # Деление группы сверяется по подписям узлов, а не по id, поэтому
        # временные id первой фазы ему не мешают — но держим рядом с рёбрами:
        # обе таблицы описывают решения владельца поверх презентации.
        apply_stage_group_split(stages, reports)

    # Внешние системы этапа, названные владельцем (process-map-vjz.5). Идёт
    # ПОСЛЕ сборки этапов: таблица дописывает в stage["inputs"], которые к этому
    # моменту уже собраны автоматикой, и проверяет, что запись не задвоилась.
    overview_texts = {
        normalize_text(shape.text).casefold()
        for shape in overview_shapes
        if shape.has_text
    }
    apply_owner_decision_external_io(stages, overview_texts, overview_report)

    # Правая колонка выходов этапа (SPEC §4.2) — блоки выходов слайда 2.
    #
    # Стоят они под боксами ГРУПП, а не под контейнером этапа (process-map-2of):
    # у этапа 1 это [182] под «Преобразование прогноза…» и [184] под «Обогащение
    # прогноза клиентскими заказами», у этапа 3 — [152], [153] и [163] под тремя
    # своими группами. Все лежат в одной горизонтальной полосе T≈4.61–4.67 млн
    # EMU и относятся к этапу целиком, поэтому на выборку это не влияет — но
    # искать их под контейнером этапа бесполезно.
    # direction='out' — по происхождению: это блок выходов этапа в обзоре
    # (см. NodeDraft.direction, задача process-map-24p). Абсцисса такого блока
    # к колонке отношения не имеет: у этапов 1 и 2 она левее середины области
    # шагов, и прежнее геометрическое правило зачисляло эти узлы во входы.
    #
    # Артефакт, который уже есть среди узлов этапа, вторым узлом не заводится
    # (проверка по casefold ниже) и направления НЕ меняет: на слайде
    # детализации он нарисован в левой колонке входов, и переписать ему
    # direction на 'out' значило бы решить за владельца процесса, что
    # презентация в одном из двух мест ошибается. Такие узлы перечислены в
    # отчёте (report.dedup_key_outputs).
    for index, stage in enumerate(stages):
        blocks = overview.output_blocks[index] if index < len(overview.output_blocks) else []
        inbound = {
            node["label"].casefold()
            for node in stage["nodes"]
            if node.get("direction") == "in"
        }
        stage["keyOutputs"] = choose_key_outputs(blocks, inbound)
        existing = {node["label"].casefold() for node in stage["nodes"]}
        added: list[dict] = []
        for paragraphs, box, sid in blocks:
            # step считается по ВСЕМ абзацам, включая заголовок: высота блока на
            # слайде поделена между ними, и пропуск заголовка не должен сдвигать
            # остальные пункты с их настоящих слайдовых координат.
            step = box.height / max(len(paragraphs), 1)
            items_start = block_items_start(paragraphs)
            for para_index, para in enumerate(paragraphs):
                if para_index < items_start:
                    # Заголовок перечня — не артефакт процесса (process-map-t9j).
                    continue
                if para.casefold() in existing:
                    twin = next(
                        (n for n in stage["nodes"] if n["label"].casefold() == para.casefold()),
                        None,
                    )
                    if twin is not None and twin.get("direction") == "in":
                        overview_report.dedup_key_outputs.append(
                            f"этап {stage['number']}: «{para}» назван выходом этапа в обзоре "
                            f"(слайд 2), но на слайде детализации это узел «{twin['id']}» из "
                            f"колонки входов — оставлен входом, отдельный узел-выход не создан"
                        )
                    continue
                existing.add(para.casefold())
                added.append(
                    serialize_node(
                        NodeDraft(
                            node_id=ids.make(para, 2, sid, para_index),
                            node_type="data",
                            label=para,
                            box=Box(
                                box.left, int(box.top + para_index * step), box.width, int(step)
                            ),
                            direction="out",
                        )
                    )
                )
        stage["nodes"].extend(added)
        overview_report.data_nodes += len(added)
        stage["warningsCount"] = count_warnings(stage)

    # Сверка списка групп: обзор (слайд 2) против слайдов детализации.
    for index, stage in enumerate(stages):
        overview_groups = overview.group_labels[index] if index < len(overview.group_labels) else []
        detail_groups = [g["label"] for g in stage["groups"]]
        for label in [g for g in overview_groups if g not in detail_groups]:
            questions.append(
                f"этап {stage['number']}: группа «{label}» есть в обзоре (слайд 2), "
                f"но не найдена на слайде детализации"
            )
        for label in [g for g in detail_groups if g not in overview_groups]:
            questions.append(
                f"этап {stage['number']}: группа «{label}» есть на слайде детализации, "
                f"но не найдена в обзоре (слайд 2)"
            )

    # Обзорные рёбра: перевод внутренних ключей в id этапов и коды систем.
    stage_key_to_id = {f"@stage{stage['number']}": stage["id"] for stage in stages}
    system_by_sid = {f"@sys:{s['code']}:{s['sid']}": s for s in overview.systems}

    resolved: list[tuple[str, str, str]] = []
    for edge in overview.edges:
        source, target = edge["source"], edge["target"]
        s_key, t_key = stage_key_to_id.get(source), stage_key_to_id.get(target)
        s_sys, t_sys = system_by_sid.get(source), system_by_sid.get(target)
        if s_key and t_key:
            resolved.append((s_key, t_key, "process"))
        elif s_key and t_sys:
            direction = register_system(stages, t_sys, source, "out")
            resolved.append(
                (s_key, t_sys["code"], "integration")
                if direction == "out"
                else (t_sys["code"], s_key, "integration")
            )
        elif s_sys and t_key:
            direction = register_system(stages, s_sys, target, "in")
            resolved.append(
                (s_sys["code"], t_key, "integration")
                if direction == "in"
                else (t_key, s_sys["code"], "integration")
            )
        else:
            overview_report.lines_skipped.append(
                f"слайд 2: ребро {source} → {target} — не пара «этап/система»"
            )

    overview_edges: list[dict] = []
    seen: set[tuple[str, str]] = set()
    for source, target, kind in resolved:
        if (source, target) in seen:
            continue
        seen.add((source, target))
        overview_edges.append(
            {"id": f"ov-{source}--{target}", "source": source, "target": target, "kind": kind}
        )
    # Ребро для внешней системы, названной владельцем (process-map-vjz.5).
    #
    # Линии на слайде для неё нет — и не может быть: её отсутствие и есть причина,
    # по которой автоматика систему не нашла. Но без ребра карточка встала бы в
    # свимлейн с подписью входа этапа 1 и стрелкой в этап 2: свимлейн держит одну
    # карточку на систему (collectSystems в overviewGraph.ts), и ERP этапов 1 и 2
    # делят её. Подпись при этом от одного этапа, а пунктир — к другому.
    for entry in OWNER_DECISION_EXTERNAL_IO:
        stage = next(s for s in stages if s["number"] == entry["stage"])
        pair = (
            (entry["system"], stage["id"])
            if entry["direction"] == "in"
            else (stage["id"], entry["system"])
        )
        if pair in seen:
            raise SystemExit(
                f"OWNER_DECISION_EXTERNAL_IO ({entry['task']}): ребро {pair[0]} → {pair[1]} "
                f"уже нашлось на слайде — презентация изменилась, и решение владельца "
                f"больше не нужно. Обновите таблицу в scripts/import-pptx.py."
            )
        seen.add(pair)
        overview_edges.append(
            {"id": f"ov-{pair[0]}--{pair[1]}", "source": pair[0], "target": pair[1], "kind": "integration"}
        )

    overview_report.edges = len(overview_edges)

    reports.insert(0, overview_report)
    for report in reports:
        questions.extend(report.questions)

    process_map = {
        "version": MAP_VERSION,
        "id": spec.map_id,
        "updatedAt": spec.updated_at,
        "title": spec.title,
        "moduleLabel": spec.module_label,
        "stages": stages,
        "overviewEdges": overview_edges,
    }
    return process_map, reports, questions, ids.counts


def apply_owner_decision_external_io(
    stages: list[dict],
    overview_texts: AbstractSet[str],
    report: SlideReport,
) -> None:
    """
    Добавляет внешние системы этапа, названные владельцем, а не выведенные из текста.

    ЗАЧЕМ ОТДЕЛЬНАЯ ТАБЛИЦА. ExternalIO собирается автоматически только когда в
    тексте фигуры нашлись И код системы, И направление. Для «Управление
    транзакционными данными» нет ни того, ни другого: код ERP владелец назвал
    отдельно, а слова «из модуля» в презентации нет. Дописать распознавание
    нельзя — оно начало бы срабатывать на чужих текстах.

    Расхождение с презентацией — остановка импорта, как в
    apply_owner_decision_edges и apply_stage_input_enrichment: если слайд
    поправят, решение владельца обязано упасть громко, а не рассосаться молча.
    """
    for entry in OWNER_DECISION_EXTERNAL_IO:
        label = entry["label"]
        if normalize_text(label).casefold() not in overview_texts:
            raise SystemExit(
                f"OWNER_DECISION_EXTERNAL_IO ({entry['task']}): на слайде обзора больше "
                f"нет текста «{label}» — презентация изменилась. Импорт остановлен, "
                f"чтобы решение владельца не потерялось молча: обновите таблицу в "
                f"scripts/import-pptx.py."
            )
        stage = next((s for s in stages if s["number"] == entry["stage"]), None)
        if stage is None:
            raise SystemExit(
                f"OWNER_DECISION_EXTERNAL_IO ({entry['task']}): этапа {entry['stage']} нет"
            )
        bucket = stage["inputs"] if entry["direction"] == "in" else stage["outputs"]
        if any(normalize_text(io["label"]).casefold() == normalize_text(label).casefold() for io in bucket):
            raise SystemExit(
                f"OWNER_DECISION_EXTERNAL_IO ({entry['task']}): «{label}» уже есть среди "
                f"внешних систем этапа {entry['stage']} — презентация изменилась и "
                f"запись больше не нужна. Обновите таблицу в scripts/import-pptx.py."
            )
        bucket.append(
            {
                "system": entry["system"],
                "label": label,
                "stage": entry["stage"],
                "direction": entry["direction"],
            }
        )
        report.owner_external_io.append(
            f"этап {entry['stage']}: {entry['system']} — «{label}» ({entry['source']})"
        )


def apply_stage_group_split(
    stages: list[dict],
    stage_reports: Sequence[SlideReport],
    splits: Sequence[dict] = STAGE_GROUP_SPLIT,
) -> None:
    """
    Выделяет часть узлов этапа в отдельную группу (STAGE_GROUP_SPLIT, process-map-028).

    Нужно там, где обзорный слайд показывает у этапа две группы, а слайд
    детализации рисует один контейнер: геометрии для деления нет, и вывести его
    нельзя — только записать решение владельца.

    Расхождение с презентацией останавливает импорт, как и в остальных таблицах
    решений: если шаг переименовали, молча потерять деление нельзя.
    """
    by_number = {stage["number"]: stage for stage in stages}
    for split in splits:
        report = stage_reports[split["stage"] - 1]
        stage = by_number.get(split["stage"])
        if stage is None:
            raise SystemExit(
                f"STAGE_GROUP_SPLIT ({split['task']}): этапа {split['stage']} нет "
                f"в презентации — решение владельца применить не к чему"
            )

        by_label = {node["label"]: node for node in stage["nodes"]}
        missing = [label for label in split["nodes"] if label not in by_label]
        if missing:
            raise SystemExit(
                f"STAGE_GROUP_SPLIT ({split['task']}): на этапе {split['stage']} нет "
                f"узлов {', '.join(repr(m) for m in missing)} — презентация изменилась. "
                f"Импорт остановлен, чтобы решение владельца не потерялось молча: "
                f"обновите список в scripts/import-pptx.py."
            )

        group_id = slugify(split["label"])
        if any(group["id"] == group_id for group in stage["groups"]):
            raise SystemExit(
                f"STAGE_GROUP_SPLIT ({split['task']}): группа «{split['label']}» уже есть "
                f"на этапе {split['stage']} — презентация изменилась и делить больше "
                f"не нужно. Обновите список в scripts/import-pptx.py."
            )
        stage["groups"].append({"id": group_id, "label": split["label"]})

        for label in split["nodes"]:
            node = by_label[label]
            was = node.get("group")
            node["group"] = group_id
            report.owner_groups.append(
                f"этап {split['stage']}: «{label[:60]}» — группа «{was}» → «{group_id}» "
                f"({split['task']}), на слайде детализации деления нет"
            )


def apply_owner_decision_edges(
    stages: list[dict],
    stage_reports: Sequence[SlideReport],
    decisions: Sequence[dict] = OWNER_DECISION_EDGES,
) -> None:
    """
    Досыпает в этапы рёбра из OWNER_DECISION_EDGES — единственные рёбра
    документа, которых нет в презентации (см. комментарий у самой константы).

    Инварианты те же, что у прочитанных со слайда рёбер, и проверяются здесь,
    а не только в src/data/schema.ts::validateIntegrity: оба конца — узлы ТОГО
    ЖЕ этапа, id ребра уникален (формат `e-{source}--{target}` общий с
    построенными по линиям, поэтому «уже есть такое ребро» и «дубль id» — одно
    и то же условие и разбирается одинаково).

    Узел, названный в решении, но исчезнувший из презентации, — ОСТАНОВКА
    импорта, а не пропуск: пропустить значило бы тихо потерять решение
    владельца, а именно ради того, чтобы оно не терялось, список и заведён.
    """
    by_number = {stage["number"]: stage for stage in stages}
    for decision in decisions:
        report = stage_reports[decision["stage"] - 1]
        stage = by_number.get(decision["stage"])
        if stage is None:
            raise SystemExit(
                f"OWNER_DECISION_EDGES ({decision['task']}): этапа {decision['stage']} нет "
                f"в презентации — решение владельца применить не к чему"
            )
        node_ids = {node["id"] for node in stage["nodes"]}
        existing = {edge["id"] for edge in stage["edges"]}
        endpoints = (decision["source"], *decision["targets"])
        missing = [node_id for node_id in endpoints if node_id not in node_ids]
        if missing:
            raise SystemExit(
                f"OWNER_DECISION_EDGES ({decision['task']}): на этапе {decision['stage']} нет "
                f"узлов {', '.join(missing)} — презентация изменилась. Импорт остановлен, чтобы "
                f"решение владельца не потерялось молча: обновите список в scripts/import-pptx.py."
            )
        for target in decision["targets"]:
            edge_id = f"e-{decision['source']}--{target}"
            if edge_id in existing:
                report.owner_edges.append(
                    f"этап {decision['stage']}: {decision['source']} → {target} — уже есть "
                    f"в презентации (стрелка нарисована), решение владельца её подтверждает"
                )
                continue
            existing.add(edge_id)
            stage["edges"].append(
                {
                    "id": edge_id,
                    "source": decision["source"],
                    "target": target,
                    "kind": decision["kind"],
                }
            )
            report.owner_edges.append(
                f"этап {decision['stage']}: {decision['source']} → {target} — ДОБАВЛЕНО "
                f"по решению владельца ({decision['task']}), в презентации стрелки нет"
            )


def register_system(stages: list[dict], system: dict, stage_key: str, direction_hint: str) -> str:
    """
    Внешняя система обзора попадает в ExternalIO этапа, к которому её ведёт стрелка.
    Направление берём из текста бокса («из модуля X» / «в модуль X»); если текст
    направления не содержит (бокс «ERP»), используем направление самой стрелки.
    """
    direction = system["direction"] or direction_hint
    stage = next((s for s in stages if f"@stage{s['number']}" == stage_key), None)
    if stage is None:
        return direction
    bucket = stage["inputs"] if direction == "in" else stage["outputs"]
    add_external_io(
        bucket,
        {
            "system": system["code"],
            "label": system["label"],
            "stage": stage["number"],
            "direction": direction,
        },
    )
    return direction


# --------------------------------------------------------------------------------------
# Отчёт-сверка
# --------------------------------------------------------------------------------------


def isolated_nodes(stage: dict) -> list[dict]:
    """
    Не-data узлы этапа без единого ребра.

    Это НЕ дефект импорта: изолированный узел означает, что соответствующей
    связи нет и в самой презентации. Связи не достраиваются (CLAUDE.md запрещает
    изобретать процесс), поэтому список печатается поимённо — по нему владелец
    процесса видит, где связь нужно проставить вручную.
    """
    connected: set[str] = set()
    for edge in stage["edges"]:
        connected.add(edge["source"])
        connected.add(edge["target"])
    return [n for n in stage["nodes"] if n["type"] != "data" and n["id"] not in connected]


def print_report(
    process_map: dict, reports: Sequence[SlideReport], questions: Sequence[str], spec: MapSpec
) -> None:
    print("=" * 78)
    print("ОТЧЁТ-СВЕРКА  scripts/import-pptx.py")
    print(f"источник: {spec.pptx.name}")
    print("=" * 78)

    for report in reports:
        print(f"\n--- слайд {report.slide_no} " + "-" * 55)
        if report.slide_no == 2:
            print(f"  этапов (контейнеров): {len(process_map['stages'])}")
            print(f"  обзорных рёбер:       {report.edges}")
            print(f"  узлов-выходов (data): {report.data_nodes}")
        else:
            print(f"  узлов:                {report.nodes} (из них data: {report.data_nodes})")
            print(f"  групп:                {report.groups}")
            print(f"  рёбер:                {report.edges} из {report.lines_total} линий/стрелок")
        if report.promoted:
            print(f"  узлы из текстбоксов ({len(report.promoted)}):")
            for item in report.promoted:
                print(f"    + {item}")
        if report.lines_skipped:
            print(f"  линии без однозначных концов ({len(report.lines_skipped)}):")
            for item in report.lines_skipped:
                print(f"    · {item}")
        if report.loose_attachments:
            print(f"  нестрогие привязки текста ({len(report.loose_attachments)}):")
            for item in report.loose_attachments:
                print(f"    · {item}")
        if report.text_skipped:
            print(f"  ПРОПУЩЕННЫЕ ФИГУРЫ С ТЕКСТОМ ({len(report.text_skipped)}):")
            for item in report.text_skipped:
                print(f"    ! {item}")
        elif report.slide_no != 2:
            print("  пропущенных фигур с текстом: нет")

    print("\n" + "=" * 78)
    print("ИТОГО ПО ДОКУМЕНТУ")
    print("=" * 78)
    total_nodes = sum(len(s["nodes"]) for s in process_map["stages"])
    by_type: Counter[str] = Counter()
    for stage in process_map["stages"]:
        for node in stage["nodes"]:
            by_type[node["type"]] += 1
    print(f"  узлов всего: {total_nodes}")
    for key in ("step", "data", "integration", "warning"):
        print(f"    {key:<12} {by_type.get(key, 0)}")
    print(f"  рёбер в этапах: {sum(len(s['edges']) for s in process_map['stages'])}")
    print(f"  обзорных рёбер: {len(process_map['overviewEdges'])}")
    print(f"  групп:          {sum(len(s['groups']) for s in process_map['stages'])}")
    for stage in process_map["stages"]:
        print(
            f"  этап {stage['number']} «{stage['shortTitle']}»: "
            f"{len(stage['nodes'])} узлов, {len(stage['edges'])} рёбер, "
            f"{len(stage['groups'])} групп, {len(stage['inputs'])} входов, "
            f"{len(stage['outputs'])} выходов, warningsCount={stage.get('warningsCount', '—')}"
        )
    # Направление data-узлов (SPEC §3, задача process-map-24p). Печатается
    # отдельным блоком, потому что это ответ на вопрос «почему на экране этапа
    # столько-то входов и столько-то выходов», а прежний ответ («так легли
    # координаты») больше не действует.
    print("\n" + "=" * 78)
    print("НАПРАВЛЕНИЕ data-УЗЛОВ (node.direction) — КОЛОНКИ ВХОДОВ И ВЫХОДОВ")
    print("=" * 78)
    print("  Ставится по происхождению фигуры, не по координатам:")
    print("    'in'  — левая колонка входов слайда детализации;")
    print("    'out' — блок выходов этапа под его контейнером на слайде обзора.")
    unset = 0
    for stage in process_map["stages"]:
        data_nodes = [n for n in stage["nodes"] if n["type"] == "data"]
        ins = sum(1 for n in data_nodes if n.get("direction") == "in")
        outs = sum(1 for n in data_nodes if n.get("direction") == "out")
        unset += len(data_nodes) - ins - outs
        flow = len(stage["nodes"]) - len(data_nodes)
        print(
            f"  этап {stage['number']} «{stage['shortTitle']}»: поток {flow}, "
            f"входов {ins}, выходов {outs}"
        )
    if unset:
        print(f"  ВНИМАНИЕ: data-узлов без direction: {unset} — это ошибка импортёра:")
        print("  все data-узлы рождаются в одном из двух мест, и оба знают направление.")
    else:
        print("  data-узлов без direction: нет")
    dedup = [item for report in reports for item in report.dedup_key_outputs]
    if dedup:
        print(f"  артефакт назван выходом в обзоре, но существует входом ({len(dedup)}):")
        for item in dedup:
            print(f"    · {item}")

    # Рёбра, которых в презентации НЕТ. Отдельный блок и есть то место, по
    # которому через полгода видно, что источник этих связей — решение
    # владельца процесса, а не стрелка на слайде (задача process-map-7bz).
    owner_edges = [item for report in reports for item in report.owner_edges]
    print("\n" + "=" * 78)
    print("РЁБРА ПО РЕШЕНИЮ ВЛАДЕЛЬЦА ПРОЦЕССА — В ПРЕЗЕНТАЦИИ ИХ НЕТ")
    print("=" * 78)
    if owner_edges:
        print("  Источник этих связей — OWNER_DECISION_EDGES в scripts/import-pptx.py,")
        print("  а НЕ линия на слайде. Всё остальное в документе прочитано из презентации.")
        for item in owner_edges:
            print(f"    · {item}")
    else:
        print("  список OWNER_DECISION_EDGES пуст — все рёбра прочитаны из презентации")

    # Входы, взятые со слайда ОБЗОРА вместо слайда детализации. Текст всё равно
    # из презентации, но выбор источника сделал владелец — блок и есть то место,
    # по которому это видно через полгода (задача process-map-qjl).
    owner_inputs = [item for report in reports for item in report.owner_inputs]
    print("\n" + "=" * 78)
    print("ВХОДЫ ПО СЛАЙДУ ОБЗОРА — НА СЛАЙДЕ ДЕТАЛИЗАЦИИ ИХ НЕТ ИЛИ ОНИ КОРОЧЕ")
    print("=" * 78)
    if owner_inputs:
        print("  Источник — STAGE_INPUT_ENRICHMENT в scripts/import-pptx.py: текст взят")
        print("  из презентации, но с другого слайда, чем остальная колонка входов.")
        for item in owner_inputs:
            print(f"    · {item}")
    else:
        print("  список STAGE_INPUT_ENRICHMENT пуст — все входы прочитаны со слайдов детализации")

    # Внешние системы, названные владельцем (process-map-vjz.5): кода системы
    # или направления в тексте презентации нет, автоматика их взять не могла.
    print("\n" + "=" * 78)
    print("ВНЕШНИЕ СИСТЕМЫ ПО РЕШЕНИЮ ВЛАДЕЛЬЦА — В ТЕКСТЕ НЕТ КОДА ИЛИ НАПРАВЛЕНИЯ")
    print("=" * 78)
    owner_external_io = [item for report in reports for item in report.owner_external_io]
    if owner_external_io:
        print("  Источник — OWNER_DECISION_EXTERNAL_IO в scripts/import-pptx.py.")
        for item in owner_external_io:
            print(f"    · {item}")
    else:
        print("  список OWNER_DECISION_EXTERNAL_IO пуст")

    # Деление группы, которого на слайде детализации нет (process-map-028).
    owner_groups = [item for report in reports for item in report.owner_groups]
    print("\n" + "=" * 78)
    print("ГРУППЫ ПО РЕШЕНИЮ ВЛАДЕЛЬЦА — НА СЛАЙДЕ ДЕТАЛИЗАЦИИ ДЕЛЕНИЯ НЕТ")
    print("=" * 78)
    if owner_groups:
        print("  Источник — STAGE_GROUP_SPLIT в scripts/import-pptx.py: обзорный слайд")
        print("  показывает две группы, слайд детализации — один контейнер.")
        for item in owner_groups:
            print(f"    · {item}")
    else:
        print("  список STAGE_GROUP_SPLIT пуст — все группы прочитаны со слайдов детализации")

    # Рёбра, разрешённые явной привязкой коннектора (process-map-3wh.16).
    cxn = [item for report in reports for item in report.cxn_edges]
    print("\n" + "=" * 78)
    print("РЁБРА ПО ПРИВЯЗКАМ КОННЕКТОРОВ — НЕ ПО ГЕОМЕТРИИ")
    print("=" * 78)
    if cxn:
        print(f"  {len(cxn)} рёбер, у которых хотя бы один конец взят из stCxn/endCxn.")
        print("  Привязка точнее геометрии: указывает на фигуру, а не на точку.")
        print("  Привязка к контейнеру группы разрешается в ближайший шаг этой группы;")
        print("  веерную связь «в группу целиком» коннектор выразить не может — для неё")
        print("  остаётся OWNER_DECISION_EDGES.")
        for item in cxn:
            print(f"    · {item}")
    else:
        print("  привязок в презентации нет — все связи выведены геометрически")

    # Узлы, ставшие интеграциями не по заливке, а по коду системы (7v1).
    promoted = [item for report in reports for item in report.promoted_integrations]
    print("\n" + "=" * 78)
    print("ИНТЕГРАЦИИ ПО КОДУ СИСТЕМЫ — ЗАЛИВКА В ПРЕЗЕНТАЦИИ ОБЫЧНАЯ")
    print("=" * 78)
    if promoted:
        print("  Серая заливка A6A6A6 есть только у входящих интеграций слайдов 3-4.")
        print("  Исходящие нарисованы как обычные шаги, поэтому тип выведен из кода системы.")
        for item in promoted:
            print(f"    · {item}")
    else:
        print("  повышений не было — все интеграции распознаны по заливке или тексту")

    print("\n" + "=" * 78)
    print("ИЗОЛИРОВАННЫЕ не-data УЗЛЫ — где в презентации связи нет")
    print("=" * 78)
    print("  Связи для них не достраиваются: в исходнике их действительно нет.")
    print("  Эти узлы нужно связать вручную или подтвердить, что связи быть не должно.")
    total_orphans = 0
    for stage in process_map["stages"]:
        orphans = isolated_nodes(stage)
        total_orphans += len(orphans)
        print(f"\n  этап {stage['number']} «{stage['shortTitle']}»: {len(orphans)}")
        for node in orphans:
            group = node.get("group")
            print(f"    · {node['id']}")
            print(
                f"        {node['type']:<11} «{node['label']}»"
                + (f"   [группа: {group}]" if group else "")
            )
    print(f"\n  всего изолированных не-data узлов: {total_orphans}")

    if questions:
        print("\n" + "=" * 78)
        print("ОТКРЫТЫЕ ВОПРОСЫ (нужны задачи -t question)")
        print("=" * 78)
        for item in questions:
            print(f"  ? {item}")


# --------------------------------------------------------------------------------------
# Перенос ручных полей из предыдущего process.json
# --------------------------------------------------------------------------------------


@dataclass
class CarryOverReport:
    """Что произошло с полями, которых в презентации нет."""

    had_previous: bool = False
    previous_nodes: int = 0
    transferred: list[str] = field(default_factory=list)
    cleared: list[str] = field(default_factory=list)
    invalid: list[str] = field(default_factory=list)
    lost: list[str] = field(default_factory=list)
    screens_transferred: int = 0
    screens_lost: int = 0


def is_screen_link(value: object) -> bool:
    """ScreenLink из schema.ts: ровно { title: string, url: string }."""
    return (
        isinstance(value, dict)
        and set(value) == {"title", "url"}
        and isinstance(value.get("title"), str)
        and isinstance(value.get("url"), str)
    )


def is_owner(value: object) -> bool:
    return isinstance(value, str)


FIELD_VALIDATORS = {"screen": is_screen_link, "owner": is_owner}


def reorder_keys(payload: dict, order: Sequence[str]) -> dict:
    """
    Пересобирает словарь в порядке ключей zod-схемы. Нужен, потому что перенос
    добавляет `owner`/`screen` уже после `position`, а порядок ключей влияет на
    байты файла (см. NODE_KEY_ORDER).
    """
    result = {key: payload[key] for key in order if key in payload}
    # Ключей вне схемы быть не должно; если появились — не теряем их молча.
    result.update({key: value for key, value in payload.items() if key not in result})
    return result


def load_previous_map(path: Path) -> dict | None:
    """
    Предыдущий process.json. Отсутствие файла — норма (первый запуск на чистом
    репозитории). А вот битый файл — НЕ норма: перезаписать его молча значит
    ровно то, ради чего эта функция написана, поэтому импорт останавливается.
    """
    if not path.exists():
        return None
    try:
        previous = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, ValueError) as error:
        raise SystemExit(
            f"Предыдущий {path} не читается ({error}). Импорт остановлен, чтобы не "
            f"затереть ручные ссылки на экраны: почините или удалите файл вручную."
        ) from error
    if not isinstance(previous, dict):
        raise SystemExit(f"Предыдущий {path} — не объект карты процесса. Импорт остановлен.")
    return previous


def index_previous(previous: dict) -> tuple[dict[str, dict], dict[str, dict]]:
    """id → узел и id → этап предыдущего документа."""
    nodes: dict[str, dict] = {}
    stages: dict[str, dict] = {}
    for stage in previous.get("stages") or []:
        if not isinstance(stage, dict):
            continue
        if isinstance(stage.get("id"), str):
            stages[stage["id"]] = stage
        for node in stage.get("nodes") or []:
            if isinstance(node, dict) and isinstance(node.get("id"), str):
                nodes[node["id"]] = node
    return nodes, stages


def carry_field(target: dict, source: dict, name: str, where: str, report: CarryOverReport) -> None:
    """
    Переносит одно поле. Три РАЗЛИЧИМЫХ состояния источника:

      · ключа нет           — поля никогда не было, переносить нечего, молчим;
      · значение null       — человек СОЗНАТЕЛЬНО удалил ссылку (SPEC §4.4,
                              кнопка «Удалить ссылку» пишет screen: null).
                              В сам JSON null не попадает: ProcessNodeSchema
                              объявляет screen как .optional(), а не .nullable(),
                              и «ссылки нет» кодируется отсутствием ключа.
                              Поэтому поле не переносим, но и не считаем потерей —
                              печатаем отдельной строкой, чтобы отличать от
                              «ссылки не было»;
      · значение есть       — переносим, если оно валидно по schema.ts.
    """
    if name not in source:
        return
    value = source[name]
    if value is None:
        report.cleared.append(f"{where}: {name} — было явно удалено (null), оставлено пустым")
        return
    if not FIELD_VALIDATORS[name](value):
        report.invalid.append(f"{where}: {name} = {value!r} — не проходит schema.ts, НЕ перенесено")
        return
    target[name] = value
    if name == "screen":
        report.screens_transferred += 1
        report.transferred.append(f"{where}: screen → «{value['title']}» {value['url']}")
    else:
        report.transferred.append(f"{where}: {name} → «{value}»")


def describe_lost(node: dict, field_name: str) -> str:
    value = node.get(field_name)
    if field_name == "screen" and isinstance(value, dict):
        return f"screen «{value.get('title')}» {value.get('url')}"
    return f"{field_name} «{value}»"


def carry_over_manual_fields(fresh: dict, previous: dict | None) -> CarryOverReport:
    """
    Переносит PRESERVED_* поля из предыдущего документа в свежесобранный.
    Сопоставление строго по `id` (узлы — глобально по документу, этапы — по
    stage.id): id стабильны по построению, а привязка по порядку обхода или по
    подписи развалилась бы при первой же правке презентации.

    Возвращает отчёт: что перенесено, что было явно очищено, что невалидно и
    что потеряно вместе с исчезнувшим узлом.
    """
    report = CarryOverReport()
    if previous is None:
        return report
    report.had_previous = True

    prev_nodes, prev_stages = index_previous(previous)
    report.previous_nodes = len(prev_nodes)

    fresh_node_ids: set[str] = set()
    fresh_labels: dict[str, list[str]] = {}
    for stage in fresh["stages"]:
        prev_stage = prev_stages.get(stage["id"])
        if prev_stage is not None:
            for name in PRESERVED_STAGE_FIELDS:
                carry_field(stage, prev_stage, name, f"этап «{stage['id']}»", report)
        for node in stage["nodes"]:
            fresh_node_ids.add(node["id"])
            fresh_labels.setdefault(node["label"], []).append(node["id"])
            prev_node = prev_nodes.get(node["id"])
            if prev_node is None:
                continue
            for name in PRESERVED_NODE_FIELDS:
                carry_field(node, prev_node, name, f"узел «{node['id']}»", report)

    # Узлы, которых в презентации больше нет. Молча потерять ссылку нельзя —
    # печатаем id, подпись и сам url, чтобы её можно было проставить заново,
    # и подсказываем узел с такой же подписью, если он появился под новым id.
    for node_id, prev_node in sorted(prev_nodes.items()):
        if node_id in fresh_node_ids:
            continue
        for name in PRESERVED_NODE_FIELDS:
            if prev_node.get(name) is None:
                continue
            if name == "screen":
                report.screens_lost += 1
            hint = ""
            twins = [i for i in fresh_labels.get(prev_node.get("label", ""), []) if i != node_id]
            if twins:
                hint = f"; возможно, это теперь {', '.join(twins)}"
            report.lost.append(
                f"узла «{node_id}» больше нет в презентации — потеряно "
                f"{describe_lost(prev_node, name)} (подпись: «{prev_node.get('label', '')}»){hint}"
            )

    for stage in fresh["stages"]:
        stage["nodes"] = [reorder_keys(node, NODE_KEY_ORDER) for node in stage["nodes"]]
    fresh["stages"] = [reorder_keys(stage, STAGE_KEY_ORDER) for stage in fresh["stages"]]
    return report


def print_carry_over(report: CarryOverReport, spec: MapSpec) -> None:
    print("\n" + "=" * 78)
    print("РУЧНЫЕ ПОЛЯ (ссылки на экраны, ответственные) — ПЕРЕНОС ИЗ ПРЕДЫДУЩЕГО JSON")
    print("=" * 78)
    if not report.had_previous:
        print(f"  предыдущего {spec.json.name} нет — первый запуск, переносить нечего")
        return
    print(f"  узлов в предыдущем файле: {report.previous_nodes}")
    print(f"  перенесено ссылок (screen): {report.screens_transferred}")
    print(f"  потеряно ссылок (screen):   {report.screens_lost}")
    if report.transferred:
        print(f"  перенесено полей ({len(report.transferred)}):")
        for item in report.transferred:
            print(f"    + {item}")
    else:
        print("  перенесённых полей нет")
    if report.cleared:
        print(f"  явно удалённые ранее ссылки ({len(report.cleared)}) — это НЕ потеря:")
        for item in report.cleared:
            print(f"    · {item}")
    if report.invalid:
        print(f"  НЕВАЛИДНЫЕ ЗНАЧЕНИЯ ({len(report.invalid)}):")
        for item in report.invalid:
            print(f"    ! {item}")
    if report.lost:
        print(f"  ПОТЕРЯННЫЕ РУЧНЫЕ ПОЛЯ ({len(report.lost)}) — проставить заново:")
        for item in report.lost:
            print(f"    ! {item}")
    else:
        print("  потерянных ручных полей нет")


def print_layout_required(process_map: dict, in_pipeline: bool) -> None:
    """
    Импорт — ПЕРВАЯ половина конвейера. Требование прогнать раскладку печатается
    последним блоком (его видно, даже если отчёт-сверку пролистали) и говорит,
    что именно сейчас лежит в файле, а не просто «не забудьте».

    `in_pipeline` — скрипт запущен из scripts/data.ts, раскладка стартует сразу
    после: пугать нечем, но сказать, что файл ещё сырой, всё равно надо.
    """
    nodes = sum(len(stage["nodes"]) for stage in process_map["stages"])
    without_slide = sum(
        1
        for stage in process_map["stages"]
        for node in stage["nodes"]
        if "slidePosition" not in node
    )
    print("\n" + "=" * 78)
    if in_pipeline:
        print("ШАГ 1 ИЗ 2 ГОТОВ — ДАЛЬШЕ РАСКЛАДКА (scripts/layout.ts, запускается сейчас)")
    else:
        print("КОНВЕЙЕР НЕ ЗАВЕРШЁН — ОБЯЗАТЕЛЬНО: npm run layout")
    print("=" * 78)
    print(f"  узлов: {nodes}; в position сейчас СЫРАЯ геометрия слайда — карточки")
    print("  на ней накладываются друг на друга и показывать её нельзя;")
    print("  пригодные координаты считает scripts/layout.ts (dagre).")
    print("  исходная геометрия сохранена в node.slidePosition, раскладка сидируется ею")
    if without_slide:
        print(f"  ВНИМАНИЕ: узлов без slidePosition: {without_slide} — это ошибка импортёра")
    if not in_pipeline:
        print("\n  одной командой:   npm run data     (import-pptx.py → layout.ts)")
        print("  сторож в тестах:  tests/mapContract.test.ts сверяет координаты,")
        print("                    так что незавершённый конвейер делает npm run check красным")


# --------------------------------------------------------------------------------------
# Запись файлов
# --------------------------------------------------------------------------------------


def write_json(path: Path, payload: object) -> None:
    text = json.dumps(payload, ensure_ascii=False, indent=2) + "\n"
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding="utf-8", newline="\n")


def check_unique_ids(process_map: dict) -> None:
    node_ids: set[str] = set()
    edge_ids: set[str] = set()
    for stage in process_map["stages"]:
        for node in stage["nodes"]:
            if node["id"] in node_ids:
                raise SystemExit(f"Неуникальный id узла: {node['id']}")
            node_ids.add(node["id"])
        for edge in stage["edges"]:
            if edge["id"] in edge_ids:
                raise SystemExit(f"Неуникальный id ребра: {edge['id']}")
            edge_ids.add(edge["id"])
    for edge in process_map["overviewEdges"]:
        if edge["id"] in edge_ids:
            raise SystemExit(f"Неуникальный id ребра: {edge['id']}")
        edge_ids.add(edge["id"])


def collect_required_node_ids(process_map: dict) -> list[str]:
    """
    SPEC §7: обязательный минимум — id ШАГОВ (step/warning/integration).
    data-узлы получаются дроблением текстбоксов-списков, их состав — эвристика,
    и фиксировать его тестом нельзя.
    """
    ids: list[str] = []
    for stage in process_map["stages"]:
        ids.extend(n["id"] for n in stage["nodes"] if n["type"] != "data")
    return ids


# --------------------------------------------------------------------------------------
# Самопроверка переноса (python scripts/import-pptx.py --self-test)
# --------------------------------------------------------------------------------------


def _fresh_fixture() -> dict:
    """Свежесобранный документ в том виде, в каком его отдаёт build_process_map."""
    return {
        "version": MAP_VERSION,
        "id": MAP_ID,
        "updatedAt": MAP_UPDATED_AT,
        "title": MAP_TITLE,
        "moduleLabel": MAP_MODULE_LABEL,
        "stages": [
            {
                "id": "stage-1",
                "number": 1,
                "title": "Этап 1",
                "shortTitle": "Этап 1",
                "keyOutputs": [],
                "warningsCount": 0,
                "groups": [],
                "nodes": [
                    {"id": "kept", "type": "step", "label": "Шаг", "position": {"x": 1, "y": 2}},
                    {"id": "cleared", "type": "step", "label": "Ш2", "position": {"x": 3, "y": 4}},
                    {"id": "never", "type": "step", "label": "Ш3", "position": {"x": 5, "y": 6}},
                    {"id": "bad", "type": "step", "label": "Ш4", "position": {"x": 7, "y": 8}},
                    {"id": "renamed-2-9", "type": "step", "label": "Ушёл", "position": {"x": 9, "y": 9}},
                ],
                "edges": [],
                "inputs": [],
                "outputs": [],
            }
        ],
        "overviewEdges": [],
    }


def _previous_fixture() -> dict:
    link = {"title": "Экран плана", "url": "https://inplan.example/plan"}
    return {
        "version": MAP_VERSION,
        "id": MAP_ID,
        "updatedAt": MAP_UPDATED_AT,
        "title": MAP_TITLE,
        "moduleLabel": MAP_MODULE_LABEL,
        "stages": [
            {
                "id": "stage-1",
                "number": 1,
                "title": "Этап 1",
                "shortTitle": "Этап 1",
                "keyOutputs": [],
                "warningsCount": 0,
                "screen": {"title": "Обзор этапа", "url": "https://inplan.example/stage-1"},
                "groups": [],
                "nodes": [
                    {
                        "id": "kept",
                        "type": "step",
                        "label": "Шаг",
                        "owner": "Планировщик спроса",
                        "screen": link,
                        "position": {"x": 0, "y": 0},
                    },
                    # null — «пользователь удалил ссылку», а не «ссылки не было».
                    {
                        "id": "cleared",
                        "type": "step",
                        "label": "Ш2",
                        "screen": None,
                        "position": {"x": 0, "y": 0},
                    },
                    # ключа screen нет вовсе.
                    {"id": "never", "type": "step", "label": "Ш3", "position": {"x": 0, "y": 0}},
                    # значение не проходит ScreenLinkSchema.
                    {
                        "id": "bad",
                        "type": "step",
                        "label": "Ш4",
                        "screen": {"url": "https://inplan.example/x"},
                        "position": {"x": 0, "y": 0},
                    },
                    # узел, которого в презентации больше нет, но ссылка была.
                    {
                        "id": "gone",
                        "type": "step",
                        "label": "Ушёл",
                        "screen": {"title": "Старый экран", "url": "https://inplan.example/old"},
                        "position": {"x": 0, "y": 0},
                    },
                ],
                "edges": [],
                "inputs": [],
                "outputs": [],
            }
        ],
        "overviewEdges": [],
    }


def run_self_test() -> int:
    """Проверки переноса ручных полей. Только stdlib, презентация не нужна."""
    checks = 0

    def check(condition: bool, message: str) -> None:
        nonlocal checks
        checks += 1
        if not condition:
            raise SystemExit(f"САМОПРОВЕРКА ПРОВАЛЕНА: {message}")

    # 1. Контракт с serialize_node: импортёр не создаёт ручных полей.
    produced = serialize_node(NodeDraft(node_id="x", node_type="step", label="L", box=Box(0, 0, 1, 1)))
    check(
        set(produced) <= set(IMPORTER_NODE_FIELDS),
        f"serialize_node отдаёт ключи вне IMPORTER_NODE_FIELDS: {set(produced) - set(IMPORTER_NODE_FIELDS)}",
    )
    check(
        not (set(PRESERVED_NODE_FIELDS) & set(IMPORTER_NODE_FIELDS)),
        "PRESERVED_NODE_FIELDS пересекается с IMPORTER_NODE_FIELDS",
    )

    # 1b. Исходная геометрия слайда сохраняется отдельным полем и совпадает с
    # position В МОМЕНТ ИМПОРТА (дальше position перезапишет npm run layout).
    box = Box(2 * EMU_PER_PX, 3 * EMU_PER_PX, 10 * EMU_PER_PX, 10 * EMU_PER_PX)
    laid = serialize_node(NodeDraft(node_id="y", node_type="step", label="L", box=box))
    check("slidePosition" in laid, "serialize_node не пишет slidePosition")
    check(
        laid["slidePosition"] == {"x": 2, "y": 3},
        f"slidePosition не равен геометрии слайда: {laid['slidePosition']}",
    )
    check(laid["position"] == laid["slidePosition"], "position при импорте ≠ slidePosition")
    # Разные объекты: раскладка меняет position и не должна задеть slidePosition.
    laid["position"]["x"] = 999
    check(laid["slidePosition"]["x"] == 2, "position и slidePosition — один и тот же объект")

    # 1c. direction (задача process-map-24p): поле пишется у data-узлов и
    # только у них, и стоит на своём месте в NODE_KEY_ORDER.
    step_node = serialize_node(NodeDraft(node_id="s", node_type="step", label="L", box=Box(0, 0, 1, 1)))
    check("direction" not in step_node, "direction проставлен не-data узлу")
    for value in ("in", "out"):
        data_node = serialize_node(
            NodeDraft(node_id="d", node_type="data", label="L", box=Box(0, 0, 1, 1), direction=value)
        )
        check(data_node.get("direction") == value, f"direction={value} не записан")
        check(
            list(data_node) == [k for k in NODE_KEY_ORDER if k in data_node],
            f"порядок ключей data-узла нарушен: {list(data_node)}",
        )
    check("direction" in NODE_KEY_ORDER, "direction отсутствует в NODE_KEY_ORDER")
    check("direction" in IMPORTER_NODE_FIELDS, "direction не объявлен полем импортёра")

    # 1d. Рёбра по решению владельца процесса (задача process-map-7bz).
    def _stage_fixture() -> list[dict]:
        return [
            {
                "number": 3,
                "nodes": [{"id": "src"}, {"id": "a"}, {"id": "b"}],
                "edges": [{"id": "e-src--a", "source": "src", "target": "a", "kind": "process"}],
            }
        ]

    # Решение-фикстура передаётся параметром, а не подменой глобали: константу
    # OWNER_DECISION_EDGES читает ещё и tests/importPreserve.test.ts (регуляркой
    # по исходнику, без Python), и второй похожий литерал в файле сбил бы разбор.
    fake_decisions = (
        {
            "task": "self-test",
            "stage": 3,
            "source": "src",
            "targets": ("a", "b"),
            "kind": "process",
            "why": "самопроверка",
        },
    )
    reports = [SlideReport(slide_no=n + 3) for n in range(4)]
    stages = _stage_fixture()
    apply_owner_decision_edges(stages, reports, fake_decisions)
    edge_ids = [e["id"] for e in stages[0]["edges"]]
    check(edge_ids == ["e-src--a", "e-src--b"], f"недостающее ребро не добавлено: {edge_ids}")
    check(len(reports[2].owner_edges) == 2, "в отчёт попали не все концы решения")
    check(
        any("уже есть в презентации" in item for item in reports[2].owner_edges),
        "совпавшее со слайдом ребро не отмечено как пришедшее из презентации",
    )
    check(
        any("ДОБАВЛЕНО" in item for item in reports[2].owner_edges),
        "добавленное ребро не отмечено как решение владельца",
    )

    # Повтор ничего не дублирует: id ребра — функция концов.
    apply_owner_decision_edges(stages, [SlideReport(slide_no=n + 3) for n in range(4)], fake_decisions)
    check(
        [e["id"] for e in stages[0]["edges"]] == edge_ids,
        "повторное применение решения продублировало рёбра",
    )

    # Исчезнувший узел — остановка импорта, а не тихий пропуск.
    broken = _stage_fixture()
    broken[0]["nodes"] = [{"id": "src"}, {"id": "a"}]
    try:
        apply_owner_decision_edges(
            broken, [SlideReport(slide_no=n + 3) for n in range(4)], fake_decisions
        )
    except SystemExit as error:
        check("b" in str(error), "в сообщении нет id пропавшего узла")
    else:
        check(False, "пропавший узел решения не остановил импорт")

    # Объявление верхнего уровня разбирается тем же способом, что и в vitest:
    # у каждого решения есть основание-задача и непустой список концов.
    for decision in OWNER_DECISION_EDGES:
        check(
            decision["task"].startswith("process-map-"),
            f"решение без задачи-основания: {decision['task']}",
        )
        check(bool(decision["targets"]), f"решение {decision['task']} без концов")

    # 1c. Входы со слайда обзора (STAGE_INPUT_ENRICHMENT, process-map-qjl).
    # Сама подстановка формулировок живёт в build_stage и проверяется реальным
    # импортом; здесь — сторожа: они обязаны ронять импорт, а не молчать.
    fake_enrichment = {
        "task": "self-test",
        "stage": 4,
        "sid": 999,
        "source": "самопроверка",
        "add": ("Новая строка",),
        "expand": (("Кратко", "Полностью"),),
        "why": "самопроверка",
    }

    def _input_column() -> list[NodeDraft]:
        return [
            NodeDraft(
                node_id="col-1",
                node_type="data",
                label="Уже есть",
                box=Box(0, 0, 100, 50),
                direction="in",
            )
        ]

    drafts = _input_column()
    enrich_report = SlideReport(slide_no=6)
    apply_input_enrichment(fake_enrichment, 6, drafts, {"Кратко"}, IdFactory({}), enrich_report)
    added = [d for d in drafts if d.label == "Новая строка"]
    check(len(added) == 1, "строка со слайда обзора не добавлена в колонку входов")
    check(added[0].direction == "in", "добавленная строка не помечена входом")
    check(added[0].node_type == "data", "добавленная строка не data-узел")
    check(
        added[0].box.top == drafts[0].box.bottom,
        "добавленная строка встала не под колонкой — раскладка сидируется slidePosition",
    )
    check(
        any("ДОБАВЛЕНО" in item for item in enrich_report.owner_inputs),
        "добавленная строка не отмечена в отчёте как пришедшая со слайда обзора",
    )

    # Переформулировка не нашла своей строки — остановка, а не тихий пропуск:
    # иначе правка презентации молча вернула бы короткую подпись.
    try:
        apply_input_enrichment(
            fake_enrichment, 6, _input_column(), set(), IdFactory({}), SlideReport(slide_no=6)
        )
    except SystemExit as error:
        check("Кратко" in str(error), "в сообщении нет ненайденной формулировки")
    else:
        check(False, "ненайденная переформулировка не остановила импорт")

    # Добавляемая строка уже есть в презентации — тоже остановка: обогащение
    # устарело, и дублировать узел нельзя.
    already = _input_column()
    already[0].label = "Новая строка"
    try:
        apply_input_enrichment(
            fake_enrichment, 6, already, {"Кратко"}, IdFactory({}), SlideReport(slide_no=6)
        )
    except SystemExit as error:
        check("Новая строка" in str(error), "в сообщении нет уже существующей строки")
    else:
        check(False, "дублирующая строка не остановила импорт")

    for entry in STAGE_INPUT_ENRICHMENT:
        check(
            entry["task"].startswith("process-map-"),
            f"обогащение без задачи-основания: {entry['task']}",
        )
        check(
            bool(entry["add"]) or bool(entry["expand"]),
            f"обогащение {entry['task']} ничего не меняет",
        )

    # 2. Первый запуск: предыдущего файла нет.
    fresh = _fresh_fixture()
    empty = carry_over_manual_fields(fresh, None)
    check(not empty.had_previous and not empty.transferred and not empty.lost, "пустой перенос")
    check(fresh == _fresh_fixture(), "перенос без предыдущего файла изменил документ")

    # 3. Основной случай.
    fresh = _fresh_fixture()
    report = carry_over_manual_fields(fresh, _previous_fixture())
    nodes = {n["id"]: n for n in fresh["stages"][0]["nodes"]}

    check(
        nodes["kept"].get("screen") == {"title": "Экран плана", "url": "https://inplan.example/plan"},
        "screen не перенесён по совпадающему id",
    )
    check(nodes["kept"].get("owner") == "Планировщик спроса", "owner не перенесён")
    # 1 ссылка узла + 1 ссылка этапа (stage.screen считается тем же счётчиком).
    check(report.screens_transferred == 2, f"screens_transferred={report.screens_transferred}, ожидалось 2")

    # 4. null отличается от отсутствия ключа.
    check("screen" not in nodes["cleared"], "screen: null попал в JSON (схема его не примет)")
    check(len(report.cleared) == 1, f"явных удалений {len(report.cleared)}, ожидалось 1")
    check("cleared" in report.cleared[0], "явное удаление не названо поимённо")
    check(
        all("«never»" not in item for item in report.cleared + report.transferred + report.lost),
        "узел без ключа screen попал в отчёт — «не было» перепутано с «удалили»",
    )
    check("screen" not in nodes["never"], "узлу без ссылки ссылка приписана")

    # 5. Невалидное значение не переносится и не молчит.
    check("screen" not in nodes["bad"], "невалидный screen перенесён в JSON")
    check(len(report.invalid) == 1, f"невалидных {len(report.invalid)}, ожидалось 1")

    # 6. Исчезнувший узел: громкая потеря, url в отчёте, подсказка по подписи.
    check(report.screens_lost == 1, f"screens_lost={report.screens_lost}, ожидалось 1")
    check(len(report.lost) == 1, f"потерь {len(report.lost)}, ожидалось 1")
    check("https://inplan.example/old" in report.lost[0], "url потерянной ссылки не напечатан")
    check("renamed-2-9" in report.lost[0], "подсказка по совпадающей подписи не выдана")

    # 7. Этап.
    check(
        fresh["stages"][0].get("screen") == {"title": "Обзор этапа", "url": "https://inplan.example/stage-1"},
        "stage.screen не перенесён",
    )

    # 8. Порядок ключей — иначе экспорт из приложения перестанет совпадать побайтово.
    check(
        list(nodes["kept"]) == [k for k in NODE_KEY_ORDER if k in nodes["kept"]],
        f"порядок ключей узла нарушен: {list(nodes['kept'])}",
    )
    check(
        list(fresh["stages"][0]) == [k for k in STAGE_KEY_ORDER if k in fresh["stages"][0]],
        f"порядок ключей этапа нарушен: {list(fresh['stages'][0])}",
    )

    # 9. Идемпотентность: перенос из уже перенесённого документа ничего не меняет.
    again = json.loads(json.dumps(fresh, ensure_ascii=False))
    carry_over_manual_fields(again, json.loads(json.dumps(fresh, ensure_ascii=False)))
    check(
        json.dumps(again, ensure_ascii=False, indent=2) == json.dumps(fresh, ensure_ascii=False, indent=2),
        "повторный перенос изменил документ — идемпотентность нарушена",
    )

    # 10. Поворот коннектора (process-map-3wh.18). Python в CI не запускается,
    #     поэтому геометрия проверяется здесь: без доворота концы повёрнутой
    #     линии вычисляются по чужим координатам, и разбор выдаёт правдоподобное,
    #     но неверное ребро.
    def probe(rot: float, head: bool = False, tail: bool = True) -> Shape:
        return Shape(
            sid=1,
            kind="line",
            box=Box(1000, 2000, 400, 200),  # центр (1200, 2100)
            paragraphs=[],
            fill=None,
            flip_h=False,
            flip_v=False,
            rot=rot,
            head_arrow=head,
            tail_arrow=tail,
        )

    flat_start, flat_end = line_endpoints(probe(0))
    check(
        flat_start == (1000.0, 2000.0) and flat_end == (1400.0, 2200.0),
        f"без поворота концы должны быть углами bbox, получено {flat_start} {flat_end}",
    )

    # 90°: угол (левый верхний) уезжает вправо-вверх относительно центра.
    turned_start, turned_end = line_endpoints(probe(90))
    check(
        all(abs(a - b) < 0.5 for a, b in zip(turned_start, (1300.0, 1900.0))),
        f"поворот 90°: начало ожидалось (1300, 1900), получено {turned_start}",
    )
    check(
        all(abs(a - b) < 0.5 for a, b in zip(turned_end, (1100.0, 2300.0))),
        f"поворот 90°: конец ожидался (1100, 2300), получено {turned_end}",
    )
    check(
        turned_start != flat_start,
        "поворот не влияет на концы — доворот не применяется",
    )

    # Поворот на 360° эквивалентен отсутствию поворота.
    full_start, full_end = line_endpoints(probe(360))
    check(
        all(abs(a - b) < 0.5 for a, b in zip(full_start, flat_start))
        and all(abs(a - b) < 0.5 for a, b in zip(full_end, flat_end)),
        "поворот 360° изменил концы",
    )

    # Стрелка на начале по-прежнему разворачивает пару и после доворота.
    rev_start, rev_end = line_endpoints(probe(90, head=True, tail=False))
    check(
        rev_start == turned_end and rev_end == turned_start,
        "headEnd не разворачивает концы повёрнутой линии",
    )

    print(f"САМОПРОВЕРКА ПРОЙДЕНА: {checks} проверок")
    return 0


def resolve_map_spec(args: Sequence[str]) -> MapSpec:
    """
    Какую карту собираем: `--map <key>`, по умолчанию snp.

    Неизвестный ключ — остановка со списком известных, а не тихий откат на
    карту по умолчанию: молчаливая пересборка не той карты затёрла бы чужой
    файл данных.
    """
    key = DEFAULT_MAP
    if "--map" in args:
        index = list(args).index("--map")
        if index + 1 >= len(args):
            raise SystemExit("--map требует значение, например: --map snp")
        key = args[index + 1]
    spec = MAPS.get(key)
    if spec is None:
        raise SystemExit(f"Неизвестная карта «{key}». Известны: {', '.join(sorted(MAPS))}")
    return spec


def main(argv: Iterable[str]) -> int:
    args = list(argv)
    # Отчёт содержит кириллицу и стрелки: на консоли с cp866/cp1251 печать иначе
    # падает с UnicodeEncodeError уже после записи файлов.
    for stream in (sys.stdout, sys.stderr):
        if hasattr(stream, "reconfigure"):
            stream.reconfigure(encoding="utf-8", errors="replace")

    if "--self-test" in args:
        return run_self_test()

    spec = resolve_map_spec(args)

    # --in-pipeline ставит scripts/data.ts (npm run data): раскладка стартует
    # сразу после импорта, и требовать её отдельно уже не надо. На сам импорт
    # флаг не влияет — только на текст финального блока.

    # Ручной слой читаем ДО сборки: если предыдущий файл битый, лучше упасть
    # раньше, чем после разбора презентации.
    previous = load_previous_map(spec.json)

    # Профиль разбора: устройство презентации у карт разное (MapSpec.profile).
    build = build_single_slide_map if spec.profile == "single-slide" else build_process_map

    # Фаза 1 — подсчёт коллизий базовых slug'ов, фаза 2 — стабильные id.
    _, _, _, collisions = build(None, spec)
    process_map, reports, questions, _ = build(dict(collisions), spec)

    carry_over = carry_over_manual_fields(process_map, previous)

    check_unique_ids(process_map)
    write_json(spec.json, process_map)
    write_json(spec.required_nodes, collect_required_node_ids(process_map))
    print_report(process_map, reports, questions, spec)
    print_carry_over(carry_over, spec)
    print(f"\nзаписано: {spec.json.relative_to(ROOT).as_posix()}")
    print(f"записано: {spec.required_nodes.relative_to(ROOT).as_posix()}")
    print_layout_required(process_map, "--in-pipeline" in args)
    if carry_over.lost:
        print(
            f"\nВНИМАНИЕ: потеряно ручных полей: {len(carry_over.lost)} "
            f"(из них ссылок на экраны: {carry_over.screens_lost}). "
            f"Список выше — проставьте их заново в редакторе. Код возврата {EXIT_LINKS_LOST}."
        )
        return EXIT_LINKS_LOST
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
